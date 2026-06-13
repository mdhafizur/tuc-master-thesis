"""Build the master-thesis defense slides as a .pptx.

Target: 20 minutes, 10 main slides @ 2 min each + a few backup slides for Q&A.
Design: clean white, TU Chemnitz green accent, no clipart.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "defense_slides.pptx"

# TU Chemnitz green (approx) and ink for body text
TUC_GREEN = RGBColor(0x00, 0x6E, 0x4F)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
RULE = RGBColor(0xCC, 0xCC, 0xCC)

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank_layout = prs.slide_layouts[6]  # blank


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, font="Calibri", align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=20,
                color=INK, font="Calibri", bold_first=False, indent_more=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        # item can be a string or a (text, level) tuple
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        bullet = ("• " if level == 0 else "– ")
        run = p.add_run()
        run.text = bullet + text
        run.font.name = font
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.bold = bold_first and level == 0
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_rule(slide, top, color=RULE, height_pt=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), top,
                                  SW - Inches(1.2), Pt(height_pt))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    return shp


def add_accent_bar(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(0.45),
                                  Inches(0.12), Inches(0.55))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = TUC_GREEN
    return shp


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.6), SH - Inches(0.45),
             Inches(7), Inches(0.3),
             "Md Hafizur Rahman · Code Guardian · TU Chemnitz · 11.05.2026",
             size=10, color=MUTED)
    add_text(slide, SW - Inches(1.6), SH - Inches(0.45),
             Inches(1.0), Inches(0.3),
             f"{page_no} / {total}",
             size=10, color=MUTED)


def title_slide(text_title, subtitle, footer_meta):
    s = prs.slides.add_slide(blank_layout)
    # accent block
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(0.8),
                              Inches(0.18), Inches(1.6))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = TUC_GREEN

    add_text(s, Inches(0.95), Inches(0.75), SW - Inches(1.5), Inches(1.0),
             "Master's Thesis Defense", size=18, color=TUC_GREEN, bold=True)
    add_text(s, Inches(0.95), Inches(1.15), SW - Inches(1.5), Inches(2.0),
             text_title, size=30, bold=True, color=INK)
    add_text(s, Inches(0.95), Inches(2.95), SW - Inches(1.5), Inches(0.6),
             subtitle, size=18, color=MUTED)
    add_rule(s, Inches(3.85))
    add_bullets(s, Inches(0.95), Inches(4.05), SW - Inches(1.5), Inches(2.5),
                footer_meta, size=16, color=INK)
    add_text(s, Inches(0.95), SH - Inches(0.55), SW - Inches(1.5), Inches(0.3),
             "Faculty of Computer Science · Chair of Distributed and "
             "Self-Organizing Systems (VSR)",
             size=10, color=MUTED)
    return s


def content_slide(title, page_no, total):
    s = prs.slides.add_slide(blank_layout)
    add_accent_bar(s)
    add_text(s, Inches(0.85), Inches(0.4), SW - Inches(1.6), Inches(0.7),
             title, size=26, bold=True, color=INK)
    add_rule(s, Inches(1.1))
    add_footer(s, page_no, total)
    return s


def set_notes(slide, notes_text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes_text


# ---------- Build the deck ----------

# Slide counts: 10 main + 4 backup = 14 total, but main footer says X / 10
TOTAL = 10

# Slide 1 — Title
s = title_slide(
    "Privacy-Preserving Source Code Vulnerability Detection and Repair "
    "using Retrieval-Augmented LLMs for Visual Studio Code",
    "Code Guardian — A local-only assistant for JavaScript / TypeScript",
    [
        "Md Hafizur Rahman · Student ID 810641 · M.Sc. Web Engineering",
        "",
        "Internal Supervisor:  Abubaker Gaber",
        "Internal Examiner:  Dr.-Ing. Sebastian Heil",
        "Chair:  Prof. Dr.-Ing. Martin Gaedke",
        "",
        "Defense · 11.05.2026",
    ],
)
set_notes(s, """Good afternoon. My name is Md Hafizur Rahman. I will defend my master's thesis on Code Guardian, a privacy-preserving secure-coding assistant for VS Code that runs entirely on the developer's machine. Total time today is 20 minutes for the talk. I'll cover the problem, the system, the evaluation setup, the headline results across the five requirements, and the limitations. Then we'll have time for questions.""")

# Slide 2 — Problem
s = content_slide("The Problem", 2, TOTAL)
add_bullets(
    s, Inches(0.85), Inches(1.4), SW - Inches(1.7), Inches(5.5),
    [
        "Cloud LLM assistants (Copilot, ChatGPT) require sending source code "
        "off the developer's machine.",
        ("Regulatory and confidentiality blockers in finance, healthcare, "
         "defence, and air-gapped environments.", 1),
        "SAST tools (Semgrep, CodeQL) are local but cannot reason across "
        "context or generate repairs.",
        ("On our JS/TS corpus: Semgrep reaches 12.68 % canonical recall.", 1),
        "No existing tool combines: local inference, retrieval-augmented "
        "reasoning, and developer-controlled repair.",
        ("This thesis builds and evaluates such a tool.", 1),
    ],
    size=18,
)
set_notes(s, """The motivation has three legs. First, cloud LLMs require sending source code off-site — that's a blocker in many regulated industries. Second, SAST tools like Semgrep and CodeQL are local but pattern-based — on our corpus Semgrep reaches only 12.68 % canonical recall. Third, no existing tool combines all three of: local inference, retrieval-augmented reasoning, and developer-controlled repair. This thesis builds and empirically evaluates exactly that combination.""")

# Slide 3 — Research questions + requirements
s = content_slide("Research Questions and Requirements", 3, TOTAL)
add_text(s, Inches(0.85), Inches(1.3), Inches(6), Inches(0.4),
         "Research Questions", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(1.7), Inches(6), Inches(4),
    [
        "RQ1  Can a local LLM on consumer hardware reach acceptable "
        "detection quality under a strict privacy envelope?",
        "RQ2  Does grounding the LLM in retrieved CWE / OWASP knowledge "
        "(RAG) improve detection?",
        "RQ3  Is in-IDE latency acceptable for developer workflows?",
    ],
    size=16,
)
add_text(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.4),
         "Five Requirements", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(7.2), Inches(1.7), Inches(5.5), Inches(4),
    [
        "R1  Detection accuracy   (precision, recall, F1)",
        "R2  Detection consistency   (JSON parse, inter-run)",
        "R3  Repair quality   (manual review + auto-applicable)",
        "R4  Usability   (latency, IDE integration)",
        "R5  Privacy   (no exfiltration, signed corpus)",
    ],
    size=16,
)
add_rule(s, Inches(5.9))
add_text(s, Inches(0.85), Inches(6.0), SW - Inches(1.7), Inches(0.9),
         "Requirements R1–R5 are operationalised against four-level threshold "
         "scales fixed in Chapter 2 before evaluation begins.",
         size=14, color=MUTED)
set_notes(s, """The thesis is organised around three research questions and five requirements. The research questions ask: (1) feasibility — can a local LLM on consumer hardware deliver useful detection? (2) grounding — does retrieval help? (3) practicality — is latency acceptable? The five requirements operationalise these questions with measurable thresholds. R1 is accuracy, R2 consistency, R3 repair, R4 usability, R5 privacy. All thresholds are fixed before evaluation — no post-hoc adjustment.""")

# Slide 4 — System architecture
s = content_slide("System Architecture", 4, TOTAL)
add_bullets(
    s, Inches(0.85), Inches(1.35), SW - Inches(1.7), Inches(5.7),
    [
        "VS Code extension orchestrates four developer workflows.",
        ("Real-time function diagnostics  ·  on-demand file scan  ·  "
         "interactive Q&A panel  ·  workspace audit.", 1),
        "Two-stage local pipeline.",
        ("Stage 1 — detection: JSON-mode LLM call, optional RAG context, "
         "consensus filter across three runs.", 1),
        ("Stage 2 — repair: separate structured call, returns "
         "{code, language}; developer applies via Quick Fix.", 1),
        "Privacy-by-construction.",
        ("Ollama bound to loopback; no telemetry; embeddings local; "
         "vulnerability data refresh fetches public metadata only.", 1),
        ("RAG corpus signed with Ed25519; container pins Node 20.19.0-alpine "
         "and locks dependencies via npm ci.", 1),
    ],
    size=17,
)
set_notes(s, """The system is a VS Code extension with four developer workflows: real-time function diagnostics, on-demand file scan, interactive Q&A, and workspace audit. Under the hood, every workflow uses the same two-stage pipeline: stage 1 detects vulnerabilities, stage 2 generates the repair on demand. Detection runs three passes under deterministic decoding and applies consensus filtering. Privacy is architectural, not bolted-on. Ollama binds to loopback, embeddings are local, the RAG corpus carries an Ed25519 signature, and the whole stack runs in a pinned container. The only outbound channel is the optional public-metadata refresh — which never sees user code.""")

# Slide 5 — Evaluation setup
s = content_slide("Evaluation Setup", 5, TOTAL)
add_bullets(
    s, Inches(0.85), Inches(1.35), SW - Inches(1.7), Inches(5.7),
    [
        "Corpus.",
        ("Curated JS / TS  ·  101 cases (71 vulnerable + 30 secure)  ·  "
         "20 CWE categories.", 1),
        ("External  ·  15 cases from OWASP NodeGoat, OWASP Juice Shop, "
         "three named CVEs.", 1),
        ("Whole-project scan against OWASP NodeGoat for end-to-end "
         "validation.", 1),
        "Configurations.",
        ("5 Ollama models  ×  {LLM-only, LLM + RAG}  =  10 configurations  ·  "
         "+ 3 SAST baselines (Semgrep, CodeQL, ESLint).", 1),
        ("3 runs per sample under deterministic decoding "
         "(temperature = 0, seed = 42)  →  303 invocations per configuration.", 1),
        "Statistical rigour.",
        ("Held-out 71 / 30 split keeps threshold tuning off the test set.", 1),
        ("McNemar with Bonferroni across paired model comparisons  ·  "
         "exact-binomial CIs for small-n rates.", 1),
    ],
    size=15,
)
set_notes(s, """The evaluation uses a curated JS/TS corpus of 101 cases — 71 vulnerable and 30 secure — across 20 CWE categories. Note: the original task description named Juliet and OWASP Benchmark; neither has a JavaScript port, so the substitution was agreed with the supervisor. Independent validation comes from 15 external cases drawn from NodeGoat, Juice Shop, and three named CVEs, plus a whole-project scan of NodeGoat. I evaluate 5 Ollama models in two modes — with and without RAG — giving 10 LLM configurations, plus three SAST baselines. Each runs three times under deterministic decoding for a total of 303 invocations per config. Statistical discipline: a held-out 71/30 split keeps threshold tuning off the test set, and I apply Bonferroni correction across paired model comparisons.""")

# Slide 6 — R1: Detection accuracy + RAG ablation
s = content_slide("Results — Detection Accuracy (R1)", 6, TOTAL)
add_text(s, Inches(0.85), Inches(1.3), Inches(6), Inches(0.4),
         "Headline:  qwen3:8b + RAG", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(1.75), Inches(6), Inches(4.5),
    [
        "F1   71.43 %    (precision 72.46 %, recall 70.42 %, FPR 10 %)",
        "With confidence gate ≥ 0.67 :  F1 74.07 %, precision 78.13 %.",
        "On the held-out 71 / 30 test split :  F1 61.11 %, FPR 0 %.",
        "Outperforms every SAST baseline on recall-driven F1.",
        ("Semgrep canonical recall :  12.68 %.", 1),
    ],
    size=15,
)
add_text(s, Inches(7.4), Inches(1.3), Inches(5.5), Inches(0.4),
         "RAG ablation across 5 models", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(7.4), Inches(1.75), Inches(5.5), Inches(4.5),
    [
        "qwen3:8b      +3.91  F1",
        "gemma3:1b   ≈ 0",
        "gemma3:4b   ≈ 0",
        "qwen3:4b      −5.22  F1",
        "codellama   −29.89  F1",
        "Only the codellama RAG-degradation survives "
        "Bonferroni correction.",
    ],
    size=15,
)
add_rule(s, Inches(6.3))
add_text(s, Inches(0.85), Inches(6.4), SW - Inches(1.7), Inches(0.8),
         "Conclusion : retrieval is not a uniform win — its benefit is "
         "model-dependent. Treat RAG as a per-model design choice, not a "
         "default.",
         size=13, color=INK, bold=True)
set_notes(s, """The headline configuration is qwen3:8b with RAG. It reaches F1 71.43 %, precision 72.46 %, recall 70.42 % at a 10 % false-positive rate. Applying a confidence gate that requires two out of three runs to agree lifts F1 to 74.07 % and precision to 78.13 %. On the frozen held-out test split, F1 drops to 61.11 % but FPR also drops to zero. Now the RAG ablation is where it gets interesting. RAG lifts qwen3:8b by about four F1 points, is neutral on the gemma3 family, slightly degrades qwen3:4b, and severely degrades codellama by nearly 30 points. Only the codellama degradation survives Bonferroni correction. The conclusion for practitioners: don't add RAG by reflex; test it per model.""")

# Slide 7 — R3: Repair quality
s = content_slide("Results — Repair Quality (R3)", 7, TOTAL)
add_text(s, Inches(0.85), Inches(1.3), SW - Inches(1.7), Inches(0.4),
         "Two complementary metrics", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(1.75), SW - Inches(1.7), Inches(5),
    [
        "Manual review  ·  generous bar  ·  n = 25",
        ("Semantic correctness :  89.5 %      (fix addresses the CWE).", 1),
        ("Combined correctness :  68 %        "
         "(semantic + strategy + executable).", 1),
        "Auto-applicable rate  ·  strict bar  ·  fully automated",
        ("90.32 %  on issued fixes   (n = 279).", 1),
        ("84.85 %  across all stage-2 calls   (n = 297, "
         "no-fix abstentions count as failures).", 1),
        "Structured contract  ·  Stage 2 returns "
        "{ code, language } and 18 explicit abstentions.",
        ("Repair is opt-in :  every applied patch is developer-confirmed via "
         "Quick Fix.", 1),
    ],
    size=16,
)
add_rule(s, Inches(6.4))
add_text(s, Inches(0.85), Inches(6.5), SW - Inches(1.7), Inches(0.7),
         "Two-tier metric so a reader can pick the bar that matches their "
         "workflow  ·  manual for decision support, auto-applicable for "
         "deployment.",
         size=13, color=MUTED, bold=True)
set_notes(s, """For repair quality I report two metrics in parallel. The manual review on 25 samples — a generous bar — gives 89.5 % semantic correctness, meaning the fix addresses the right CWE. The stricter combined bar — semantic plus strategy plus executable code — drops to 68 %. The auto-applicable rate is fully automated, no human in the loop: 90.32 % on issued fixes, or 84.85 % across all stage-2 calls when 18 explicit no-fix abstentions count as failures. The two-tier approach lets a reader pick the bar that matches their use case. Important: every applied patch is developer-confirmed through VS Code's Quick Fix UI — the system never edits code autonomously.""")

# Slide 8 — R4 latency + R5 privacy
s = content_slide("Results — Usability (R4) and Privacy (R5)", 8, TOTAL)
add_text(s, Inches(0.85), Inches(1.3), Inches(6), Inches(0.4),
         "R4 — Stage-1 detection latency", size=16, bold=True,
         color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(1.75), Inches(6), Inches(4.5),
    [
        "Headline qwen3:8b + RAG  ·  median 2,216 ms  ·  p95 4,327 ms.",
        "On-demand band ≤ 5 s  ·  comfortably inside.",
        "Interactive band ≤ 1.5 s  ·  gemma3:1b + RAG (1,019 ms) and "
        "qwen3:4b (1,328 ms) qualify.",
        "Real-time band ≤ 500 ms  ·  SAST baselines only.",
        "Stage 2 repair adds ~1.5 × the stage-1 cost on demand.",
    ],
    size=14,
)
add_text(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.4),
         "R5 — Privacy", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(7.2), Inches(1.75), Inches(5.5), Inches(4.5),
    [
        "Loopback-only Ollama binding   (network isolation arm of threat "
        "model).",
        "Ed25519-signed RAG corpus manifest   (provenance arm).",
        "Prompt-injection harness   12 cases   ·   leakFreeRate 100 %.",
        "Zero non-loopback transmission across the full corpus run.",
        "Container pinned to Node 20.19.0-alpine   ·   "
        "npm ci   ·   seed = 42.",
    ],
    size=14,
)
add_rule(s, Inches(6.4))
add_text(s, Inches(0.85), Inches(6.5), SW - Inches(1.7), Inches(0.7),
         "Resource footprint :  harness CPU 3 ms median, RSS 87 MB  ·  "
         "Ollama dominates at 6 – 8 GB RAM for an 8B-q4 model.",
         size=13, color=MUTED)
set_notes(s, """Latency. Headline configuration sits at 2,216 ms median for stage-1 detection — comfortably within the on-demand band of five seconds. The interactive sub-1.5-second band is reachable by gemma3:1b plus RAG and by qwen3:4b, at a detection-quality cost. Real-time sub-500-ms feedback remains the exclusive territory of SAST baselines. Privacy is operationalised across four checks. Ollama is bound to loopback only. The RAG corpus is signed with Ed25519. A 12-case prompt-injection harness gives 100 % leakFreeRate. And across the entire corpus run, no non-loopback transmission was observed. The harness side of the system is lightweight — three milliseconds of CPU, 87 MB of RSS. Ollama itself dominates the resource budget at six to eight gigabytes of RAM for an 8B model at four-bit quantization.""")

# Slide 9 — Contributions, limitations, future work
s = content_slide("Contributions  ·  Limitations  ·  Future Work", 9, TOTAL)
add_text(s, Inches(0.85), Inches(1.3), Inches(4.2), Inches(0.4),
         "Contributions", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(1.7), Inches(4.2), Inches(5),
    [
        "Working VS Code extension  ·  local LLM + RAG + signed corpus + "
        "container.",
        "Two-stage pipeline with structured JSON contracts.",
        "Empirical study across 5 models × 2 modes + 3 SAST baselines on a "
        "CWE-mapped JS / TS corpus.",
        "Reproducibility infrastructure  ·  byte-identical runs, signed "
        "manifest, pinned container.",
        "Statistical discipline  ·  Bonferroni, exact-binomial CIs, "
        "held-out split.",
    ],
    size=13,
)
add_text(s, Inches(5.25), Inches(1.3), Inches(3.9), Inches(0.4),
         "Limitations", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(5.25), Inches(1.7), Inches(3.9), Inches(5),
    [
        "Manual repair review by single reviewer   (n = 25).",
        "Curated corpus  ·  selection bias mitigated by held-out split and "
        "external 15-case set.",
        "One real-world project   (NodeGoat).",
        "Consensus filter is an artefact under deterministic decoding  ·  "
        "needs seed rotation.",
        "No cloud-LLM baseline  ·  excluded by privacy threat model.",
    ],
    size=13,
)
add_text(s, Inches(9.35), Inches(1.3), Inches(3.4), Inches(0.4),
         "Future Work", size=16, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(9.35), Inches(1.7), Inches(3.4), Inches(5),
    [
        "Seed rotation for genuine consensus signal.",
        "Inter-rater agreement on manual review.",
        "Larger real-world validation across projects.",
        "Streaming partial results to break the real-time band.",
        "Adaptive top-k retrieval and per-model RAG gating.",
    ],
    size=13,
)
set_notes(s, """Five contributions: a working extension, the two-stage pipeline with structured contracts, the empirical study across 10 LLM configurations and three baselines, reproducibility infrastructure with byte-identical runs and a signed manifest, and the statistical discipline including Bonferroni correction. Five limitations to be honest about. The 25-sample manual review used a single reviewer — that is the methodology's weakest link. The curated corpus carries selection-bias risk; we mitigate with the held-out split and the external 15-case set, but it remains a single curator. Only one real-world project. The consensus filter contributes no discriminative signal under deterministic decoding — that is acknowledged in the thesis. And there is no cloud-LLM baseline, by design — the privacy threat model excludes it. Future work flows directly from these limitations: seed rotation to make the consensus filter meaningful, inter-rater agreement, more real-world projects, and architectural ideas like streaming and per-model RAG gating.""")

# Slide 10 — Questions
s = content_slide("Thank you  ·  Questions?", 10, TOTAL)
add_text(s, Inches(0.85), Inches(2.6), SW - Inches(1.7), Inches(1.2),
         "Code Guardian", size=40, bold=True, color=TUC_GREEN)
add_text(s, Inches(0.85), Inches(3.5), SW - Inches(1.7), Inches(0.8),
         "Privacy-preserving secure-coding assistance, on the developer's "
         "machine.",
         size=20, color=INK)
add_rule(s, Inches(4.6))
add_text(s, Inches(0.85), Inches(4.8), SW - Inches(1.7), Inches(0.6),
         "Headline numbers", size=14, bold=True, color=TUC_GREEN)
add_bullets(
    s, Inches(0.85), Inches(5.15), SW - Inches(1.7), Inches(2),
    [
        "F1  71.43 %   ·   repair  89.5 % semantic / 90.32 % auto-applicable   "
        "·   latency  2.2 s median   ·   leakFreeRate  100 %",
        "Bonferroni-significant cross-model effect :  RAG degrades codellama "
        "by 29.89 F1 points.",
    ],
    size=14,
)
set_notes(s, """Thank you. I'm happy to take questions on any aspect — the substituted corpus, the single-reviewer manual review, the codellama RAG effect, the choice not to compare against cloud LLMs, the resource budget, or anything else.""")

# ---------- Backup slides for Q&A ----------

def backup(title):
    s = prs.slides.add_slide(blank_layout)
    add_accent_bar(s)
    add_text(s, Inches(0.85), Inches(0.4), SW - Inches(1.6), Inches(0.7),
             f"Backup  ·  {title}", size=22, bold=True, color=TUC_GREEN)
    add_rule(s, Inches(1.1))
    return s


# B1 — single-reviewer rebuttal (the biggest examiner attack)
s = backup("Why the single-reviewer manual review is bounded, not broken")
add_bullets(
    s, Inches(0.85), Inches(1.4), SW - Inches(1.7), Inches(5.5),
    [
        "Manual review is the qualitative complement, not the headline.",
        ("Headline R3 metric is auto-applicable rate "
         "(90.32 %, n = 279)  ·  fully objective, fully automated.", 1),
        "Rubric is documented and replicable.",
        ("Four dimensions  ·  semantic correctness, strategy alignment, "
         "executability, conciseness.", 1),
        ("All 25 reviewed cases plus rubric are in the reproduction "
         "package.", 1),
        "Scope-bounded by master's-thesis resources.",
        ("Inter-rater agreement on a 10-sample subset is the obvious "
         "next step  ·  flagged in future work.", 1),
        "The 25-sample number is reported with explicit limitation language "
        "in Chapter 5 and the conclusion.",
    ],
    size=15,
)

# B2 — why curated corpus
s = backup("Why a curated corpus rather than an existing JS benchmark")
add_bullets(
    s, Inches(0.85), Inches(1.4), SW - Inches(1.7), Inches(5.5),
    [
        "Original task corpora (Juliet, OWASP Benchmark) are C / C++ / Java "
        "only  ·  no first-party JS port.",
        "SecurityEval and similar JS sets are small and partially synthetic; "
        "their CWE labelling is inconsistent with OWASP Top 10 grouping.",
        "Selection-bias mitigations in this thesis :",
        ("Held-out 71 / 30 split with a frozen TEST set keeps threshold "
         "tuning out of the headline numbers.", 1),
        ("15-case external corpus from NodeGoat, Juice Shop, "
         "and three named CVEs  ·  independent of the curator.", 1),
        ("Whole-project NodeGoat run gives a project-level recall figure "
         "independent of sample selection.", 1),
        ("Corpus provenance (URL + commit) is bound into the Ed25519-signed "
         "manifest  ·  cannot be silently rewritten.", 1),
    ],
    size=15,
)

# B3 — no cloud comparison
s = backup("Why no cloud-LLM baseline")
add_bullets(
    s, Inches(0.85), Inches(1.4), SW - Inches(1.7), Inches(5.5),
    [
        "Cloud-LLM confidentiality is the problem statement  ·  introducing "
        "a cloud baseline contradicts the threat model.",
        "The deterministic floor is provided by three SAST baselines  ·  "
        "Semgrep, CodeQL, ESLint.",
        "Code Guardian wins recall-driven F1 versus every SAST baseline on "
        "the same canonical-taxonomy matcher.",
        "Frontier cloud models will always win raw F1  ·  the thesis "
        "competes on the constraint set (local, zero-egress, "
        "reproducible).",
        "Adding a published-numbers comparison against GPT-4 on similar "
        "corpora would be a one-paragraph addendum  ·  doable post-defense.",
    ],
    size=15,
)

# B4 — consensus / R2 artefact
s = backup("Why R2 reports 100 % inter-run agreement")
add_bullets(
    s, Inches(0.85), Inches(1.4), SW - Inches(1.7), Inches(5.5),
    [
        "Decoding is locked  ·  temperature = 0, seed = 42, "
        "format = JSON, fixed prompt.",
        "Under these settings the three runs are byte-identical  ·  the "
        "consensus filter therefore contributes no discriminative "
        "signal in this corpus.",
        "This is acknowledged explicitly in evaluation.tex and in the "
        "appendix reproducibility section.",
        "Consensus filtering becomes meaningful under seed rotation, which "
        "is left to future work.",
        "Practical implication  ·  the confidence-gate band currently "
        "filters by output-presence consistency rather than model "
        "self-consistency.",
    ],
    size=15,
)

prs.save(OUT)
print(f"Wrote {OUT}")
