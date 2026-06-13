"""Builds defense_slides_v3.pptx — Tier 1+2 redesign over v2.

Changes vs v2:
- Real chart images on slides 15, 16, 18 (instead of placeholders)
- Real architecture diagram on slide 11
- Real threat-model diagram on slide A5
- Per-model results table (A1) populated from canonical evaluation results
- Confusion matrices (A2) with accurate case-level counts
- Slide 17 (R3 Repair) and 19 (R5 Privacy) restructured as Huzefa-style
  metric bars
- New References slide between Contributions and Thank You
- Footer stripped to "Md Hafizur Rahman" + page number (no title/uni/date noise)
- Slide 2 / 10 placeholders kept with clearer drop-in instructions
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os
import re

# ---- design tokens ----
NAVY = RGBColor(0x14, 0x21, 0x3D)
NAVY_DARK = RGBColor(0x0B, 0x14, 0x28)
ACCENT = RGBColor(0xFC, 0xA3, 0x11)
RED = RGBColor(0xC9, 0x18, 0x4A)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
# bright status colours (legible on the navy hero cards)
STATUS_GREEN = RGBColor(0x3F, 0xB9, 0x50)
STATUS_AMBER = RGBColor(0xF5, 0xB7, 0x14)
STATUS_RED = RGBColor(0xFF, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xFB)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE6, 0xEE)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
TEXT_MUTED = RGBColor(0x5B, 0x65, 0x7A)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
HEADER_H = Inches(0.9)
FOOTER_H = Inches(0.35)
MARGIN = Inches(0.5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

DATE_STR = "16.06.2026"
DATE_LONG = "June 16, 2026"
AUTHOR_SHORT = "Md Hafizur Rahman"

ASSETS = os.path.join(os.path.dirname(__file__), "assets")
ICONS = os.path.join(ASSETS, "icons")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def line_dash(line, val="dash"):
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": val}))


def line_arrow(line, w="med", length="med"):
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": w, "len": length}))


def add_rect(slide, x, y, w, h, fill=CARD_BG, line=None, line_pt=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        try:
            s.adjustments[0] = 0.06
        except Exception:
            pass
    set_solid(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, line_pt)
    s.shadow.inherit = False
    return s


def card_bg(slide, x, y, w, h, *, accent, strip=Inches(0.08),
            fill=CARD_BG, border=CARD_BORDER):
    """Rounded card with a coloured top band whose corners follow the card's
    rounded corners: an accent rounded-rect behind a white body shifted down."""
    a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        a.adjustments[0] = 0.06
    except Exception:
        pass
    set_solid(a, accent)
    no_line(a)
    a.shadow.inherit = False
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + strip, w, h - strip)
    try:
        b.adjustments[0] = 0.06
    except Exception:
        pass
    set_solid(b, fill)
    if border is None:
        no_line(b)
    else:
        set_line(b, border, 0.75)
    b.shadow.inherit = False
    return b


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
             italic=False, space_after=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_slide_chrome(slide, title, *, page_idx=None, header=True, footer=True,
                     show_title=True, accent_strip=True, show_brand=True):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
    if header:
        add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill=NAVY)
        add_text(slide, Inches(0.4), Inches(0.22), Inches(6), Inches(0.5),
                 "TECHNICAL UNIVERSITY OF CHEMNITZ", size=10, bold=True,
                 color=TEXT_LIGHT)
        if show_brand:
            add_text(slide, SLIDE_W - Inches(3.0) - Inches(0.4), Inches(0.22),
                     Inches(3.0), Inches(0.5), "CODE GUARDIAN",
                     size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    if show_title:
        title_y = HEADER_H + Inches(0.15)
        add_text(slide, MARGIN, title_y, SLIDE_W - 2 * MARGIN, Inches(0.65),
                 title, size=30, bold=True, color=NAVY)
        if accent_strip:
            add_rect(slide, MARGIN, title_y + Inches(0.65),
                     Inches(0.7), Inches(0.04), fill=ACCENT)
    if footer:
        # quiet footer: just author left, page right
        add_text(slide, MARGIN, SLIDE_H - FOOTER_H,
                 Inches(5), FOOTER_H, AUTHOR_SHORT, size=9,
                 color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        if page_idx is not None:
            add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - FOOTER_H,
                     Inches(0.8), FOOTER_H, f"{page_idx} / __",
                     size=9, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.RIGHT)


def content_box():
    y = HEADER_H + Inches(1.05)
    return (MARGIN, y, SLIDE_W - 2 * MARGIN, SLIDE_H - y - FOOTER_H - Inches(0.1))


def card(slide, x, y, w, h, *, icon=None, icon_variant="navy", title=None, body=None,
         badge=None, accent_color=NAVY, body_size=12, title_size=16,
         body_color=TEXT_DARK, title_h=Inches(0.45), body_anchor=MSO_ANCHOR.TOP,
         body_space=None):
    card_bg(slide, x, y, w, h, accent=accent_color)
    pad = Inches(0.22)
    cy = y + Inches(0.20)
    title_x = x + pad
    title_w = w - 2 * pad
    if icon:
        isize = Inches(0.42)
        iy = cy + (title_h - isize) / 2  # vertically centred against the title row
        place_icon(slide, title_x, iy, isize, icon, variant=icon_variant)
        title_x += Inches(0.58); title_w -= Inches(0.58)
    if badge:
        add_text(slide, x + w - pad - Inches(1.0), cy + Inches(0.05),
                 Inches(1.0), Inches(0.35), badge, size=10, bold=True,
                 color=accent_color, align=PP_ALIGN.RIGHT)
        title_w -= Inches(1.0)
    if title:
        add_text(slide, title_x, cy, title_w, title_h,
                 title, size=title_size, bold=True, color=NAVY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
    if body:
        by = cy + title_h + Inches(0.12)
        add_text(slide, x + pad, by, w - 2 * pad, h - (by - y) - Inches(0.2),
                 body, size=body_size, color=body_color, anchor=body_anchor,
                 space_after=body_space)


def hero_number_card(slide, x, y, w, h, *, label, value, subtext,
                     status=None, band=None):
    add_rect(slide, x, y, w, h, fill=NAVY, line=None, rounded=True)
    # status accent strip down the left edge
    if status is not None:
        add_rect(slide, x, y, Inches(0.1), h, fill=status)
    add_text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), Inches(0.4),
             label, size=12, bold=True, color=ACCENT)
    # band chip (e.g. HIGH / MEDIUM), top-right, coloured by status
    if status is not None and band is not None:
        chip_w = Inches(1.15)
        add_rect(slide, x + w - chip_w - Inches(0.22), y + Inches(0.16),
                 chip_w, Inches(0.34), fill=status, rounded=True)
        add_text(slide, x + w - chip_w - Inches(0.22), y + Inches(0.16),
                 chip_w, Inches(0.34), band, size=10, bold=True, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # big number, vertically centred in the space between the label and the
    # bottom-pinned description so the trio reads as one balanced block
    add_text(slide, x + Inches(0.25), y + Inches(0.6), w - Inches(0.5), h - Inches(1.55),
             value, size=36, bold=True, color=(status or TEXT_LIGHT),
             anchor=MSO_ANCHOR.MIDDLE)
    # description pinned to the bottom of the card
    add_text(slide, x + Inches(0.25), y + h - Inches(1.05), w - Inches(0.5), Inches(0.85),
             subtext, size=10, color=RGBColor(0xC8, 0xD0, 0xE0),
             anchor=MSO_ANCHOR.BOTTOM)


def placeholder(slide, x, y, w, h, label, *, sub=None):
    add_rect(slide, x, y, w, h, fill=RGBColor(0xEE, 0xF0, 0xF6),
             line=CARD_BORDER, line_pt=1.0, rounded=True)
    add_text(slide, x, y + h / 2 - Inches(0.35), w, Inches(0.4),
             label, size=14, bold=True,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        add_text(slide, x, y + h / 2, w, Inches(0.4),
                 sub, size=10, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def badge(slide, x, y, w, h, text, *, color=ACCENT):
    add_rect(slide, x, y, w, h, fill=color, line=None, rounded=True)
    add_text(slide, x, y, w, h, text, size=11, bold=True,
             color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_picture(slide, x, y, w, h, path):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        placeholder(slide, x, y, w, h, f"missing {os.path.basename(path)}")


def place_icon(slide, x, y, size, name, variant="navy"):
    """Place a square icon from the consistent custom set (assets/icons)."""
    path = os.path.join(ICONS, f"{name}-{variant}.png")
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=size, height=size)


def status_icon(slide, x, y, size, kind):
    """kind: 'check' | 'partial' | 'cross' — coloured matrix marker."""
    path = os.path.join(ICONS, f"st-{kind}.png")
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=size, height=size)


_MATRIX_KIND = {"✅": "check", "🟡": "partial", "❌": "cross"}


def cell_marker(slide, cx, ry, cwid, rh, token):
    kind = _MATRIX_KIND.get(token)
    if kind:
        sz = Inches(0.30)
        status_icon(slide, cx + (cwid - sz) / 2, ry + (rh - sz) / 2, sz, kind)
    else:
        add_text(slide, cx, ry, cwid, rh, token, size=16, color=NAVY_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def legend_rows(slide, lx, ly, wdt, size=12):
    rows = [("check", "High / Satisfied"), ("partial", "Partial"),
            ("cross", "Not satisfied")]
    sz = Inches(0.26)
    for i, (kind, lbl) in enumerate(rows):
        ry = ly + i * Inches(0.38)
        status_icon(slide, lx, ry, sz, kind)
        add_text(slide, lx + Inches(0.38), ry - Inches(0.03), wdt - Inches(0.4),
                 Inches(0.32), lbl, size=size, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)


def bar_metric(slide, x, y, w, h, *, label, value, max_label="", color=NAVY,
               sub=""):
    """Huzefa-style horizontal bar metric. value is a 0..1 fill ratio."""
    fill_w = int(w * max(0.05, min(1.0, value)))
    add_rect(slide, x, y, w, h, fill=LIGHT_BG, line=CARD_BORDER, rounded=True)
    add_rect(slide, x, y, fill_w, h, fill=color, line=None, rounded=True)
    # value label inside the filled bar (white text)
    add_text(slide, x + Inches(0.18), y, fill_w - Inches(0.36), h,
             label, size=11, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    # max_label outside the filled bar (muted text, right-aligned in the unfilled trailing region)
    if max_label and fill_w < w - Inches(0.4):
        add_text(slide, x + fill_w + Inches(0.1), y, w - fill_w - Inches(0.2), h,
                 max_label, size=9, color=TEXT_MUTED,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        add_text(slide, x, y + h, w, Inches(0.25),
                 sub, size=9, color=TEXT_MUTED)


# ============================================================
# SLIDE 1 — Title
# ============================================================
def slide_01_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(0.6), Inches(0.7), Inches(0.04), Inches(0.5), fill=ACCENT)
    add_text(s, Inches(0.85), Inches(0.65), Inches(8), Inches(0.5),
             "TECHNICAL UNIVERSITY OF CHEMNITZ", size=12, bold=True,
             color=TEXT_LIGHT)
    add_text(s, SLIDE_W - Inches(3.8), Inches(0.55), Inches(3.2), Inches(0.7),
             "CODE GUARDIAN", size=20, bold=True, color=TEXT_LIGHT,
             align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.85), Inches(2.0), SLIDE_W - Inches(1.7), Inches(2.0),
             "Privacy-Preserving Source Code\nVulnerability Detection & Repair",
             size=44, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(4.0), SLIDE_W - Inches(1.7), Inches(0.8),
             "Retrieval-Augmented LLMs for Visual Studio Code",
             size=20, color=RGBColor(0xC8, 0xD0, 0xE0))
    col_y = Inches(5.5)
    col_w = Inches(3.9)
    muted = RGBColor(0xC8, 0xD0, 0xE0)

    def label(x, yy, text):
        add_text(s, x, yy, col_w, Inches(0.28), text, size=9, bold=True, color=ACCENT)

    def value(x, yy, text, primary=True):
        add_text(s, x, yy, col_w, Inches(0.32), text,
                 size=12 if primary else 10, bold=primary, italic=True,
                 color=TEXT_LIGHT if primary else muted)

    x0 = Inches(0.85)
    x1 = x0 + (col_w + Inches(0.15))
    x2 = x0 + 2 * (col_w + Inches(0.15))

    # Column 1 — presenter
    label(x0, col_y, "PRESENTED BY")
    value(x0, col_y + Inches(0.30), "Md Hafizur Rahman")
    value(x0, col_y + Inches(0.62), "Mat. No. 810641  ·  M.Sc. Web Engineering", primary=False)

    # Column 2 — supervisor
    label(x1, col_y, "SUPERVISOR")
    value(x1, col_y + Inches(0.30), "Abubaker Gaber  (Internal)")

    # Column 3 — examiner + date stacked
    label(x2, col_y, "EXAMINER")
    value(x2, col_y + Inches(0.30), "Dr.-Ing. Sebastian Heil")
    label(x2, col_y + Inches(0.70), "DATE")
    value(x2, col_y + Inches(1.00), DATE_LONG)


# ============================================================
# SLIDE 2 — Problem framing (illustration placeholder kept)
# ============================================================
def slide_02_problem_illustration():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "The Developer's Dilemma", page_idx=2)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2
    ch = (h - gap) / 2

    # three problems (red) arranged TL, TR, BL
    problems = [
        ("cloud", "Can't use the cloud",
         "Sensitive code can't be pasted into a cloud model without breaking compliance."),
        ("search", "Scanners miss bugs",
         "The local scanner keeps missing real vulnerabilities, so the bug reaches production."),
        ("wrench", "Warnings, no fixes",
         "When a scanner flags an issue it only warns. Writing the repair is still your job."),
    ]
    for (ic, t, b), (c, r) in zip(problems, [(0, 0), (1, 0), (0, 1)]):
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=ic, icon_variant="red", title=t, body=b,
             body_size=13, title_size=16, accent_color=RED,
             title_h=Inches(0.5), body_anchor=MSO_ANCHOR.MIDDLE)

    # the question (bottom-right) — the hook, visually distinct in green
    qx = x + (cw + gap)
    qy = y + (ch + gap)
    add_rect(s, qx, qy, cw, ch, fill=GREEN, rounded=True)
    place_icon(s, qx + Inches(0.25), qy + Inches(0.24), Inches(0.44), "compass",
               variant="light")
    add_text(s, qx + Inches(0.85), qy + Inches(0.22), cw - Inches(1.1), Inches(0.5),
             "What's the alternative?", size=20, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, qx + Inches(0.25), qy + Inches(0.9), cw - Inches(0.5),
             ch - Inches(1.15),
             "No tool today combines privacy, full coverage, and one-click repair. "
             "Closing that gap is what this thesis sets out to do.",
             size=14, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 3 — Two failures (privacy vs coverage)
# ============================================================
def slide_03_two_failures():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Today You Pick: Send Code Away, or Get No Help", page_idx=3)
    x, y, w, h = content_box()
    gap = Inches(0.4)
    cw = (w - gap) / 2
    cards_y = y + Inches(0.45)
    card_h = Inches(3.0)
    card(s, x, cards_y, cw, card_h,
         icon="cloud", icon_variant="red", title="Cloud LLM (Copilot, ChatGPT, Cursor)",
         body=("Source code crosses the network on every request.\n\n"
               "Blocked outright in finance, healthcare, defence,\n"
               "and air-gapped environments.\n\n"
               "Strong reasoning · zero privacy guarantees."),
         body_size=14, accent_color=RED, body_anchor=MSO_ANCHOR.MIDDLE)
    card(s, x + cw + gap, cards_y, cw, card_h,
         icon="search", icon_variant="red", title="Local SAST (Semgrep, CodeQL, ESLint)",
         body=("Stays local. Cannot reason across files or generate fixes.\n\n"
               "Semgrep canonical recall on our JS/TS corpus:  12.68 %\n\n"
               "Strong privacy · poor coverage · no repair."),
         body_size=14, accent_color=RED, body_anchor=MSO_ANCHOR.MIDDLE)
    co_y = cards_y + card_h + Inches(0.45)
    add_rect(s, x, co_y, w, Inches(1.0), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), co_y, w - Inches(0.6), Inches(1.0),
             "Either choice forces a trade-off: privacy or capability, never both.",
             size=18, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — Three-way gap
# ============================================================
def slide_04_three_gap():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "No Tool Combines All Three", page_idx=4)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - 2 * gap) / 3
    rows = [
        ("lock", "Local Inference",
         "Code never leaves the\ndeveloper's machine.\nVerifiable network boundary."),
        ("book", "Retrieval-Augmented",
         "Grounded in CWE / OWASP\nknowledge.\nNot just pattern matching."),
        ("wrench", "Developer-Controlled Repair",
         "Structured fix returned.\nDeveloper applies via Quick Fix.\nNo silent rewrites."),
    ]
    cards_y = y + Inches(0.7)
    card_h = Inches(2.5)
    for i, (icon, t, b) in enumerate(rows):
        cx = x + i * (cw + gap)
        card(s, cx, cards_y, cw, card_h,
             icon=icon, title=t, body=b, body_size=14, title_size=17,
             title_h=Inches(0.85), accent_color=ACCENT,
             body_anchor=MSO_ANCHOR.MIDDLE)
    co_y = cards_y + card_h + Inches(0.45)
    add_rect(s, x, co_y, w, Inches(1.0), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), co_y, w - Inches(0.6), Inches(1.0),
             "This thesis builds and evaluates a tool that delivers all three.",
             size=18, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — RQs + Requirements
# ============================================================
def slide_05_rqs_reqs():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Research Questions and Requirements", page_idx=5)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    lw = w * 0.40 - gap / 2
    rw = w * 0.60 - gap / 2

    add_rect(s, x, y, lw, Inches(3.3), fill=CARD_BG, line=CARD_BORDER,
             line_pt=0.75, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(0.15), lw - Inches(0.6), Inches(0.4),
             "RESEARCH QUESTIONS", size=11, bold=True, color=ACCENT)
    rqs = [
        ("RQ1", "Can a local LLM on consumer hardware reach acceptable detection quality under a strict privacy envelope?"),
        ("RQ2", "Does grounding the LLM in retrieved CWE / OWASP knowledge (RAG) improve detection?"),
        ("RQ3", "Is in-IDE latency acceptable for developer workflows?"),
    ]
    row_y = y + Inches(0.65)
    for tag, q in rqs:
        add_text(s, x + Inches(0.3), row_y, Inches(0.7), Inches(0.4),
                 tag, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(1.05), row_y, lw - Inches(1.35), Inches(1.0),
                 q, size=11, color=TEXT_DARK)
        row_y += Inches(0.95)

    add_rect(s, x + lw + gap, y, rw, Inches(3.3), fill=CARD_BG,
             line=CARD_BORDER, line_pt=0.75, rounded=True)
    add_text(s, x + lw + gap + Inches(0.3), y + Inches(0.15), rw - Inches(0.6),
             Inches(0.4), "FIVE REQUIREMENTS", size=11, bold=True, color=ACCENT)
    rs = [
        ("R1", "Accuracy", "precision, recall, F1"),
        ("R2", "Consistency", "JSON parse, inter-run"),
        ("R3", "Repair Quality", "manual + auto-applicable"),
        ("R4", "Usability", "latency, IDE integration"),
        ("R5", "Privacy", "zero exfiltration, signed corpus"),
    ]
    rw_inner = (rw - Inches(0.6) - Inches(0.3)) / 2
    rx = x + lw + gap + Inches(0.3)
    ry = y + Inches(0.65)
    for i, (tag, lbl, sub) in enumerate(rs):
        col = i % 2
        row = i // 2
        cx = rx + col * (rw_inner + Inches(0.3))
        cy = ry + row * Inches(0.85)
        add_rect(s, cx, cy, rw_inner, Inches(0.75),
                 fill=LIGHT_BG, line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.15), cy + Inches(0.06), Inches(0.5), Inches(0.35),
                 tag, size=12, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.75), cy + Inches(0.06),
                 rw_inner - Inches(0.85), Inches(0.35),
                 lbl, size=13, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.75), cy + Inches(0.36),
                 rw_inner - Inches(0.85), Inches(0.35),
                 sub, size=9, color=TEXT_MUTED)

    add_rect(s, x, y + Inches(3.45), w, Inches(0.55), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(3.45), w - Inches(0.6), Inches(0.55),
             "R1–R5 map to fixed threshold scales, set before evaluation, with no post-hoc bands.",
             size=12, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    gy = y + Inches(4.15)
    add_rect(s, x, gy, w, Inches(0.85), fill=LIGHT_BG, line=CARD_BORDER, rounded=True)
    add_text(s, x + Inches(0.25), gy + Inches(0.1), Inches(2.0), Inches(0.7),
             "", size=10, bold=True, color=ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.9), gy + Inches(0.08), w - Inches(2.1), Inches(0.74),
             "RAG = look up known CWE / OWASP weakness patterns before answering   ·   "
             "CWE = standard catalogue of software-weakness types\n"
             "F1 = balance of catching real bugs vs. false alarms (higher = better)   ·   "
             "FPR = share of safe code wrongly flagged (lower = better)",
             size=10.5, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE, space_after=3)


# ============================================================
# SLIDE 6 — Existing approaches (column cards)
# ============================================================
def slide_06_existing_approaches():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Three Existing Approaches", page_idx=6)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 2 * gap) / 3
    cols = [
        ("CLOUD LLM ASSISTANTS",
         [("GitHub Copilot", "Microsoft · 2021"),
          ("ChatGPT / GPT-4", "OpenAI · 2022"),
          ("Cursor", "Anysphere · 2023"),
          ("Claude Code", "Anthropic · 2024")]),
        ("LOCAL SAST",
         [("Semgrep", "Rule-based · OSS"),
          ("CodeQL", "GitHub · taint analysis"),
          ("ESLint security plugins", "Lint-level"),
          ("Bandit / SpotBugs", "Other languages")]),
        ("ACADEMIC LLM-FOR-SEC",
         [("IRIS + GPT-4", "Li et al. · ICLR 2025"),
          ("SecRepair", "Islam et al. · NDSS 2024"),
          ("LLMSecGuard", "Kavian et al. · EASE 2024"),
          ("RESCUE", "Shi & Zhang · 2025")]),
    ]
    for i, (hdr, items) in enumerate(cols):
        cx = x + i * (cw + gap)
        add_rect(s, cx, y, cw, Inches(0.55), fill=NAVY, rounded=True)
        add_text(s, cx, y, cw, Inches(0.55), hdr, size=11, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rowy = y + Inches(0.75)
        for name, sub in items:
            add_rect(s, cx, rowy, cw, Inches(0.7), fill=CARD_BG,
                     line=CARD_BORDER, rounded=True)
            add_text(s, cx + Inches(0.2), rowy + Inches(0.08), cw - Inches(0.4),
                     Inches(0.3), name, size=12, bold=True, color=NAVY_DARK)
            add_text(s, cx + Inches(0.2), rowy + Inches(0.38), cw - Inches(0.4),
                     Inches(0.3), sub, size=10, color=TEXT_MUTED)
            rowy += Inches(0.8)


# ============================================================
# SLIDE 7 — Critical gap matrix
# ============================================================
def slide_07_gap_matrix():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "No Tool Satisfies R3 + R4 + R5 Together", page_idx=7)
    x, y, w, h = content_box()

    rows = [
        ("GitHub Copilot",   ["✅","✅","✅","🟡","❌"]),
        ("Semgrep",          ["🟡","✅","❌","✅","✅"]),
        ("CodeQL",           ["🟡","✅","❌","🟡","✅"]),
        ("IRIS (GPT-4)",     ["🟡","🟡","❌","❌","❌"]),
        ("SecRepair",        ["🟡","🟡","🟡","❌","🟡"]),
    ]
    header = ["Solution / Approach", "R1\nACCURACY", "R2\nCONSISTENCY",
              "R3\nREPAIR", "R4\nUSABILITY", "R5\nPRIVACY"]

    table_w = w * 0.75
    legend_w = w - table_w - Inches(0.3)
    col_w = [Inches(2.3)] + [(table_w - Inches(2.3)) / 5] * 5
    # size the table to fill the content area, leaving a strip for the source note
    header_h = Inches(0.75)
    table_h = h - Inches(0.55)
    row_h = (table_h - header_h) / len(rows)
    table_bottom = y + table_h

    cx = x
    for i, hd in enumerate(header):
        add_rect(s, cx, y, col_w[i], header_h, fill=NAVY)
        add_text(s, cx, y, col_w[i], header_h, hd,
                 size=11, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    ry = y + header_h
    for r_i, (name, cells) in enumerate(rows):
        is_self = name == "Code Guardian"
        bg = RGBColor(0xE9, 0xF3, 0xEA) if is_self else (LIGHT_BG if r_i % 2 == 0 else CARD_BG)
        cx = x
        add_rect(s, cx, ry, col_w[0], row_h, fill=bg, line=CARD_BORDER)
        add_text(s, cx + Inches(0.15), ry, col_w[0] - Inches(0.2), row_h,
                 name, size=13, bold=is_self, color=NAVY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[0]
        for c_i, cell in enumerate(cells):
            add_rect(s, cx, ry, col_w[c_i + 1], row_h, fill=bg, line=CARD_BORDER)
            cell_marker(s, cx, ry, col_w[c_i + 1], row_h, cell)
            cx += col_w[c_i + 1]
        ry += row_h

    lx = x + table_w + Inches(0.3)
    ly = y
    legend_h = Inches(1.7)
    add_rect(s, lx, ly, legend_w, legend_h,
             fill=CARD_BG, line=CARD_BORDER, rounded=True)
    add_text(s, lx + Inches(0.2), ly + Inches(0.15), legend_w - Inches(0.4),
             Inches(0.3), "LEGEND", size=10, bold=True, color=ACCENT)
    legend_rows(s, lx + Inches(0.2), ly + Inches(0.55), legend_w - Inches(0.4), size=12)

    co_y = ly + legend_h + Inches(0.3)
    co_h = table_bottom - co_y
    add_rect(s, lx, co_y, legend_w, co_h, fill=NAVY, rounded=True)
    add_text(s, lx + Inches(0.25), co_y + Inches(0.25), legend_w - Inches(0.5),
             Inches(0.4), "THE GAP", size=12, bold=True, color=ACCENT)
    add_text(s, lx + Inches(0.25), co_y + Inches(0.75), legend_w - Inches(0.5),
             co_h - Inches(0.95),
             "No existing tool delivers\nrepair + IDE-grade latency\n+ privacy together.",
             size=15, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # source note: where the level thresholds come from
    add_text(s, x, table_bottom + Inches(0.08), w, Inches(0.4),
             "Levels apply the fixed R1–R5 scales, set before evaluation:  "
             "accuracy / repair thresholds [Johnson 2013; Monperrus 2014],  "
             "latency [Card 1983; Nielsen 1994],  privacy = checklist.   "
             "Full scale + sources in Appendix A8.",
             size=9, color=TEXT_MUTED, italic=True)


# ============================================================
# SLIDE 8 — Four workflows on one pipeline
# ============================================================
def slide_08_workflows():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Proposed Workflow", page_idx=8)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 3 * gap) / 4
    rows = [
        ("bolt", "Real-time Diagnostics", "Inline squiggles per function as you save."),
        ("folder", "On-demand File Scan", "Right-click → scan the current file."),
        ("chat", "Interactive Q&A", "Ask about a finding in a side panel."),
        ("layers", "Workspace Audit", "Batch scan the whole project."),
    ]
    cards_y = y + Inches(0.5)
    for i, (icon, t, b) in enumerate(rows):
        cx = x + i * (cw + gap)
        card(s, cx, cards_y, cw, Inches(1.95),
             icon=icon, title=t, body=b, body_size=13, title_size=14,
             accent_color=ACCENT, body_anchor=MSO_ANCHOR.MIDDLE)

    py = cards_y + Inches(2.35)
    pipe_h = Inches(1.7)
    add_rect(s, x, py, w, pipe_h, fill=CARD_BG, line=CARD_BORDER, rounded=True)
    add_text(s, x + Inches(0.3), py + Inches(0.12), w - Inches(0.6), Inches(0.3),
             "ALL FOUR WORKFLOWS HIT THE SAME PIPELINE", size=10, bold=True,
             color=ACCENT)
    nodes = ["INPUT", "STAGE 1\nDetection (LLM + RAG)", "Consensus\nfilter",
             "STAGE 2\nRepair (LLM)", "QUICK\nFIX"]
    bx = x + Inches(0.4)
    by = py + Inches(0.55)
    bw = (w - Inches(0.8) - 4 * Inches(0.25)) / 5
    bh = Inches(0.85)
    for i, n in enumerate(nodes):
        cx = bx + i * (bw + Inches(0.25))
        fill = NAVY if i in (1, 3) else RGBColor(0xE6, 0xEB, 0xF5)
        col = TEXT_LIGHT if i in (1, 3) else NAVY
        add_rect(s, cx, by, bw, bh, fill=fill, rounded=True)
        add_text(s, cx, by, bw, bh, n, size=11, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(nodes) - 1:
            ax = cx + bw + Inches(0.02)
            add_text(s, ax, by, Inches(0.21), bh, "→", size=18,
                     color=ACCENT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 9 — Two-stage pipeline detail
# ============================================================
def slide_09_two_stage():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Two-Stage Local Pipeline", page_idx=9)
    x, y, w, h = content_box()
    gap = Inches(0.6)
    cw = (w - gap) / 2
    cards_y = y + Inches(0.45)
    card_h = Inches(2.95)
    card(s, x, cards_y, cw, card_h,
         icon="search", title="Stage 1 · Detection",
         body=("• JSON-mode LLM call  ·  schema-pinned output\n"
               "• Optional RAG context  ·  top-k CWE / OWASP snippets\n"
               "• Two consensus passes  ·  distinct seeds (42, 137)\n"
               "• Finding kept only when both passes agree\n"
               "• Returns:  category, severity, line range, evidence"),
         body_size=13, accent_color=NAVY, body_anchor=MSO_ANCHOR.MIDDLE,
         body_space=10)
    card(s, x + cw + gap, cards_y, cw, card_h,
         icon="wrench", title="Stage 2 · Repair",
         body=("• Separate structured call  ·  invoked only when detected\n"
               "• Returns  { code, language }  ·  no-fix return allowed (18 / 297)\n"
               "• Validator parses output with @babel/parser  ·  rejects prose\n"
               "• Surfaced as a VS Code Quick Fix\n"
               "• Developer applies the patch  ·  no silent rewrites"),
         body_size=13, accent_color=ACCENT, body_anchor=MSO_ANCHOR.MIDDLE,
         body_space=10)
    # pipeline arrow: detection feeds repair
    add_text(s, x + cw, cards_y, gap, card_h, "→", size=28, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    by = cards_y + card_h + Inches(0.4)
    add_rect(s, x, by, w, Inches(0.85), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by, w - Inches(0.6), Inches(0.85),
             "Privacy by construction  ·  Ollama bound to loopback  ·  Embeddings local  ·  No telemetry",
             size=14, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — VS Code screenshot
# ============================================================
def slide_10_screenshot():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Live in VS Code", page_idx=10)
    x, y, w, h = content_box()
    ss_path = os.path.join(ASSETS, "screenshot_vscode.png")
    if os.path.exists(ss_path):
        # read the real dimensions so any screenshot fits without distortion,
        # then fit inside the box left of the QR column, vertically centred
        from PIL import Image
        iw, ih = Image.open(ss_path).size
        ar = iw / ih
        avail_w = int(w - Inches(2.4))
        avail_h = int(h - Inches(0.5))
        ss_w = min(avail_w, int(avail_h * ar))
        ss_h = int(ss_w / ar)
        iy = y + (avail_h - ss_h) // 2
        add_picture(s, x, iy, ss_w, ss_h, ss_path)
    else:
        placeholder(s, x, y, w - Inches(2.2), h - Inches(0.2),
                    "Drop assets/screenshot_vscode.png here",
                    sub="Open the extension in dev mode, trigger a finding on a vulnerable JS/TS file, "
                        "capture full window with side panel + Quick Fix visible (PNG, ~1600×900)")
    qx = x + w - Inches(2.0)
    qr_path = os.path.join(ASSETS, "qr_demo.png")
    if os.path.exists(qr_path):
        add_picture(s, qx, y, Inches(2.0), Inches(2.0), qr_path)
    else:
        placeholder(s, qx, y, Inches(2.0), Inches(2.0), "QR CODE",
                    sub="Demo video / GitHub release")
    add_text(s, qx, y + Inches(2.1), Inches(2.0), Inches(0.4),
             "SCAN  ·  Install from VS Code Marketplace", size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — System architecture (real image now)
# ============================================================
def slide_11_architecture():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "System Architecture", page_idx=11)
    draw_architecture(s)


# ---- native (editable) architecture diagram --------------------------------
# Ported 1:1 from make_assets.diagram_architecture (200x80 user-unit canvas).
STORE_FILL = RGBColor(0xE9, 0xEE, 0xF8)
CONTAINER_FILL = RGBColor(0xEE, 0xF2, 0xFA)
CYL_TOP_FILL = RGBColor(0xF4, 0xF7, 0xFD)
BOUND_FILL = RGBColor(0xFB, 0xFC, 0xFE)
CREAM = RGBColor(0xFF, 0xF6, 0xE2)
SUB_ON_NAVY = RGBColor(0xC7, 0xD0, 0xE4)


def draw_architecture(slide):
    AW, AH = 200.0, 80.0                      # matplotlib canvas units
    x0, y0, fw, fh = content_box()
    fh = fh - Inches(0.2)
    # native font scale: the raster was a 15in-wide figure shown at fw inches
    SCALE = (fw / 914400.0) / 15.0
    TS = 0.88                                  # same global text shrink as raster

    def MX(v): return int(round(x0 + v / AW * fw))
    def MY(v): return int(round(y0 + (1 - v / AH) * fh))   # y is bottom-up
    def MW(v): return int(round(v / AW * fw))
    def MH(v): return int(round(v / AH * fh))
    def FS(base): return Pt(base * TS * SCALE)             # font size in points

    def _dash(line, val="dash"):
        ln = line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": val}))

    def _arrowhead(line, w="med", length="med"):
        ln = line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": w, "len": length}))

    def shape(kind, cx, cy, w, h, fill, line=None, line_pt=1.5, dash=False,
              radius=0.10):
        left, top, ww, hh = MX(cx - w / 2), MY(cy + h / 2), MW(w), MH(h)
        sp = slide.shapes.add_shape(kind, left, top, ww, hh)
        if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                sp.adjustments[0] = radius
            except Exception:
                pass
        set_solid(sp, fill)
        if line is None:
            no_line(sp)
        else:
            set_line(sp, line, line_pt)
            if dash:
                _dash(sp.line)
        sp.shadow.inherit = False
        return sp

    def settext(sp, title, sub=None, tcol=TEXT_LIGHT, scol=SUB_ON_NAVY,
                tsize=12.5, ssize=9.0):
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.03)
        tf.margin_top = tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.bold = True; r.font.size = FS(tsize)
        r.font.color.rgb = tcol; r.font.name = FONT_BODY
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(1)
            r2 = p2.add_run(); r2.text = sub
            r2.font.size = FS(ssize); r2.font.color.rgb = scol
            r2.font.name = FONT_BODY

    def box(cx, cy, w, h, title, sub=None, fill=NAVY, tcol=TEXT_LIGHT,
            scol=SUB_ON_NAVY, line=None, line_pt=1.5, tsize=12.5, ssize=9.0):
        sp = shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, w, h,
                   fill, line if line else fill, line_pt)
        settext(sp, title, sub, tcol, scol, tsize, ssize)
        return dict(cx=cx, cy=cy, L=cx - w / 2, R=cx + w / 2,
                    T=cy + h / 2, B=cy - h / 2)

    def cyl(cx, cy, w, h, title, sub):
        sp = shape(MSO_SHAPE.CAN, cx, cy, w, h, STORE_FILL, NAVY, 1.5)
        settext(sp, title, sub, tcol=NAVY, scol=TEXT_MUTED,
                tsize=10.5, ssize=8.0)
        return dict(cx=cx, cy=cy, L=cx - w / 2, R=cx + w / 2,
                    T=cy + h / 2, B=cy - h / 2)

    def conn(p0, p1, color=NAVY, width_pt=1.6, dash=False, arrow=True):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       MX(p0[0]), MY(p0[1]), MX(p1[0]), MY(p1[1]))
        c.line.color.rgb = color
        c.line.width = Pt(width_pt)
        if dash:
            _dash(c.line)
        if arrow:
            _arrowhead(c.line)
        c.shadow.inherit = False
        return c

    def floattext(cx, cy, text, color=TEXT_DARK, size=9.0, bold=True,
                  italic=False, mask=True, align=PP_ALIGN.CENTER, anchor="c"):
        cw = Pt((size * TS * SCALE) * 0.60 * len(text) + 6)
        ch = Pt((size * TS * SCALE) * 1.55)
        if align == PP_ALIGN.CENTER:
            left = MX(cx) - int(cw // 2)
        elif align == PP_ALIGN.RIGHT:
            left = MX(cx) - int(cw)
        else:
            left = MX(cx)
        tb = slide.shapes.add_textbox(left, MY(cy) - int(ch // 2), int(cw), int(ch))
        tf = tb.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = FS(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT_BODY
        if mask:
            set_solid(tb, CARD_BG)
        else:
            tb.fill.background()
        tb.line.fill.background()
        tb.shadow.inherit = False
        return tb

    def lbl(cx, cy, text, color=TEXT_DARK, size=9.0):
        return floattext(cx, cy, text, color=color, size=size, bold=True)

    # ---- boundary (drawn first, behind everything) ----
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, 4 + 178 / 2, 4.5 + 62.5 / 2, 178, 62.5,
          BOUND_FILL, NAVY, 2.0, dash=True, radius=0.04)

    OFF = 1.6
    # geometry constants (cx) shared by columns
    C_CLIENT, C_CORE, C_ANALYZER, C_OLLAMA = 24, 63, 111, 159
    # ---- connectors first so the boxes overlay their endpoints ----
    conn((37, 53 + OFF), (50, 53 + OFF), width_pt=1.9)                 # analyze
    conn((50, 53 - OFF), (37, 53 - OFF), width_pt=1.3, dash=True)
    conn((76, 53), (87, 53), width_pt=1.9)                            # on miss
    conn((135, 53 + OFF), (145, 53 + OFF), width_pt=1.9)              # prompt
    conn((145, 53 - OFF), (135, 53 - OFF), width_pt=1.3, dash=True)
    conn((61, 47), (61, 38.5), width_pt=1.9)                          # check
    conn((65, 38.5), (65, 47), width_pt=1.3, dash=True)               # hit
    conn((100, 42.75), (100, 40.5), arrow=False)                     # store elbow
    conn((100, 40.5), (73, 40.5), arrow=False)
    conn((73, 40.5), (73, 38.5))
    conn((111, 42.75), (111, 38.5), width_pt=1.9)                    # RAG enrich
    conn((111, 27.5), (111, 19.5), width_pt=1.9)                     # top-k 3
    conn((144, 14), (127, 14), width_pt=1.5, dash=True)              # sync
    conn((159, 19.5), (159, 27.5), width_pt=1.9)                    # egress req
    conn((173, 33), (176, 33), color=RED, width_pt=1.7, dash=True,  # egress out
         arrow=False)
    conn((176, 33), (176, 68.5), color=RED, width_pt=1.7, dash=True)

    # ---- analyzer container + stage cards ----
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, C_ANALYZER, 53, 48, 20.5,
          CONTAINER_FILL, NAVY, 1.6, radius=0.07)
    floattext(C_ANALYZER, 60.85, "Analyzer (two-stage)", color=NAVY_DARK,
              size=10.5, bold=True, mask=False)
    box(C_ANALYZER - 11.5, 49.5, 20, 11, "Stage 1 · Detect", "consensus · JSON",
        tsize=8.8, ssize=7.8)
    box(C_ANALYZER + 11.5, 49.5, 20, 11, "Stage 2 · Repair", "Babel-validated",
        tsize=8.8, ssize=7.8)
    conn((109.5, 49.5), (112.5, 49.5), color=ACCENT, width_pt=1.7)

    # ---- primary boxes ----
    box(C_CLIENT, 53, 26, 12, "Client UI", "VS Code · 4 workflows")
    box(C_CORE, 53, 26, 12, "Extension Core", "orchestration")
    box(C_OLLAMA, 53, 28, 12, "Ollama", "local · codellama / qwen3", ssize=8.6)
    box(C_ANALYZER, 33, 30, 11, "RAG Manager", "context retrieval",
        tsize=11.5, ssize=8.6)
    box(C_OLLAMA, 33, 28, 11, "Egress Gate", "loopback-only · off by default",
        tsize=11.0, ssize=7.6)

    # ---- cylinders ----
    cyl(C_CORE, 33, 30, 11, "Analysis Cache", "LRU · 100 · 30 min")
    cyl(C_ANALYZER, 14, 32, 11, "Knowledge Base", "HNSWlib · Ed25519-signed")
    cyl(C_OLLAMA, 14, 30, 11, "Vulnerability Data", "24 h cache · 7 sources")

    # ---- external feed (outside boundary) ----
    box(176, 73, 42, 9, "Public Security Feeds", "NVD · OWASP · CWE · GitHub",
        fill=CARD_BG, tcol=NAVY_DARK, scol=TEXT_MUTED, line=NAVY, line_pt=1.8,
        tsize=11.5, ssize=8.6)
    floattext(176, 78.9, "outside the privacy boundary", color=TEXT_MUTED,
              size=8.0, bold=False, italic=True, mask=False)

    # ---- boundary labels ----
    floattext(7, 64.6, "CODE GUARDIAN", color=NAVY_DARK, size=12.0, bold=True,
              mask=False, align=PP_ALIGN.LEFT)
    floattext(7, 61.7, "runs entirely on your laptop", color=TEXT_MUTED,
              size=9.0, bold=False, italic=True, mask=False, align=PP_ALIGN.LEFT)

    # ---- privacy callout ----
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8 + 38 / 2, 5.5 + 23 / 2, 38, 23,
          CREAM, ACCENT, 1.8, radius=0.06)
    floattext(27, 25.5, "Privacy by construction", color=NAVY_DARK, size=11.0,
              bold=True, mask=False)
    checks = ["all local execution", "no code leaves machine", "no telemetry",
              "egress off by default", "signed corpus"]
    ctb = slide.shapes.add_textbox(MX(11), MY(23.2), MW(34), MH(17))
    ctf = ctb.text_frame; ctf.word_wrap = False
    ctf.margin_left = ctf.margin_right = Pt(1)
    ctf.margin_top = ctf.margin_bottom = Pt(0)
    for i, c in enumerate(checks):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(3.0)
        rc = p.add_run(); rc.text = "✓ "
        rc.font.size = FS(9.4); rc.font.bold = True
        rc.font.color.rgb = GREEN; rc.font.name = FONT_BODY
        rt = p.add_run(); rt.text = c
        rt.font.size = FS(9.4); rt.font.color.rgb = TEXT_DARK
        rt.font.name = FONT_BODY
    ctb.fill.background(); ctb.line.fill.background()

    # ---- connector labels (on top, white-masked) ----
    lbl(43.5, 57.2, "analyze", color=NAVY)
    lbl(81.5, 55.6, "on miss", color=NAVY, size=8.6)
    lbl(140, 57.3, "prompt (JSON)", color=NAVY, size=8.0)
    lbl(54.5, 42.75, "check", color=NAVY)
    lbl(71, 42.75, "hit", color=TEXT_MUTED)
    lbl(86.5, 41.9, "store on miss", color=NAVY, size=8.6)
    lbl(122, 40.6, "RAG enrich", color=NAVY, size=8.6)
    lbl(120, 23.5, "top-k 3", color=NAVY, size=8.6)
    lbl(135.5, 16.3, "sync", color=TEXT_MUTED)
    lbl(148, 23.5, "egress req", color=NAVY, size=8.6)
    floattext(173.6, 44.2, "optional fetch", color=RED, size=8.8, bold=True,
              align=PP_ALIGN.RIGHT, mask=False)
    floattext(173.6, 41.0, "off by default", color=RED, size=8.2, bold=False,
              align=PP_ALIGN.RIGHT, mask=False)

    # ---- legend ----
    conn((60, 2.4), (67, 2.4), color=NAVY, width_pt=2.2)
    floattext(68.5, 2.4, "Local data flow", color=TEXT_DARK, size=9.4,
              bold=False, mask=False, align=PP_ALIGN.LEFT)
    conn((112, 2.4), (119, 2.4), color=RED, width_pt=1.8, dash=True)
    floattext(120.5, 2.4, "Optional outbound: public metadata, off by default",
              color=TEXT_DARK, size=9.4, bold=False, mask=False,
              align=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 12 — Privacy by Construction
# ============================================================
def slide_12_privacy_cards():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Privacy by Construction", page_idx=12)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(4.0) - gap) / 2
    items = [
        ("lock", "Loopback-only Ollama",
         "Bound to 127.0.0.1. Network isolation arm of the threat model, enforced by network policy."),
        ("seal", "Signed RAG Corpus",
         "Ed25519 manifest verified at load. Provenance arm of the threat model. The corpus cannot be silently rewritten."),
        ("shield", "Prompt-Injection Harness",
         "12 cases of crafted prompts attempting exfiltration. leakFreeRate = 100 %."),
        ("cube", "Pinned Container",
         "Node 20.19.0-alpine · npm ci · seed = 42 · byte-identical runs."),
    ]
    for i, (icon, t, b) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=icon, title=t, body=b,
             body_size=12, title_size=15, accent_color=ACCENT)


# ============================================================
# SLIDE 13 — Evaluation setup
# ============================================================
def slide_13_eval_setup():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Evaluation Setup", page_idx=13)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 2 * gap) / 3
    items = [
        ("book", "Corpus",
         "Curated JS / TS  ·  101 cases\n(71 vulnerable + 30 secure)\n20 CWE categories\n\nExternal:  15 cases from NodeGoat,\nJuice Shop, 3 named CVEs\n\nWhole-project scan:  OWASP NodeGoat"),
        ("flask", "Configurations",
         "5 Ollama models × {LLM-only, LLM+RAG}\n= 10 configurations\n\n+ 3 SAST baselines\n(Semgrep, CodeQL, ESLint)\n\n3 runs per sample  ·  deterministic\n→  303 invocations per config"),
        ("ruler", "Statistical Discipline",
         "Held-out 71 / 30 test split\n(threshold tuning off the test set)\n\nMcNemar with Bonferroni\nacross paired model comparisons\n\nExact-binomial CIs for small-n rates"),
    ]
    cards_y = y + Inches(0.5)
    for i, (icon, t, b) in enumerate(items):
        cx = x + i * (cw + gap)
        card(s, cx, cards_y, cw, Inches(3.85),
             icon=icon, title=t, body=b, body_size=11.5, title_size=15,
             accent_color=NAVY, body_anchor=MSO_ANCHOR.MIDDLE, body_space=6)


# ============================================================
# SLIDE 14 — Hero numbers (kept; numbers verified against canonical)
# ============================================================
def slide_14_hero():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Results", page_idx=14)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2
    # cards take the full remaining height minus the bottom hint strip
    bottom_strip_h = Inches(0.45)
    avail = h - Inches(0.05) - bottom_strip_h - Inches(0.15)
    ch = (avail - gap) / 2
    cards = [
        ("R1 · DETECTION",  "71.43 %",
         "F1 · qwen3:8b + RAG  ·  P 72.46 % · R 70.42 %\nOn unseen code:  held-out 61.11 %  ·  external 63.64 %",
         STATUS_AMBER, "MEDIUM"),
        ("R3 · REPAIR",     "90.32 %",
         "Auto-applicable: parses & applies (n = 279)\nSemantic-correct ≈ 68 % combined · 42 % run as-is (n = 25)",
         STATUS_GREEN, "HIGH"),
        ("R4 · LATENCY",    "2.2 s",
         "Median stage-1  ·  p95 = 4.3 s\nComfortably inside on-demand band (≤ 5 s)",
         STATUS_GREEN, "HIGH"),
        ("R5 · PRIVACY",    "100 %",
         "leakFreeRate · 12-case prompt-injection harness\nschemaValidRate 83 % (10/12) · loopback-only by policy",
         STATUS_GREEN, "MET"),
    ]
    for i, (lbl, val, sub, status, band) in enumerate(cards):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        hero_number_card(s, cx, cy, cw, ch, label=lbl, value=val, subtext=sub,
                         status=status, band=band)
    by = y + 2 * ch + gap + Inches(0.1)
    add_rect(s, x, by, w, bottom_strip_h, fill=NAVY, rounded=True)
    chip_w = Inches(1.0)
    add_rect(s, x + Inches(0.18), by + (bottom_strip_h - Inches(0.26)) / 2,
             chip_w, Inches(0.26), fill=STATUS_GREEN, rounded=True)
    add_text(s, x + Inches(0.18), by + (bottom_strip_h - Inches(0.26)) / 2,
             chip_w, Inches(0.26), "R2 · HIGH", size=9, bold=True, color=NAVY_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.3), by + Inches(0.02), w - Inches(1.6), bottom_strip_h,
             "Consistency:  100 % JSON-parse success, 3,030 invocations (inter-run agreement is deterministic by seed).   "
             "Colour = requirement band (A8).   R1 / R3 / R4 / R5 detailed next.",
             size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 — R1 detection (with real chart)
# ============================================================
def slide_15_r1_accuracy():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R1 · Detection Accuracy", page_idx=15, show_brand=False)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R1", color=ACCENT)
    x, y, w, h = content_box()
    # top: numbers card
    add_rect(s, x, y, w, Inches(1.4), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(0.15), w - Inches(0.6), Inches(0.4),
             "qwen3:8b + RAG  ·  canonical taxonomy  ·  full 101-case run",
             size=11, bold=True, color=ACCENT)
    stats = [
        ("71.43 %", "F1"),
        ("72.46 %", "Precision"),
        ("70.42 %", "Recall"),
        ("10.0 %",  "FPR"),
    ]
    sw = (w - Inches(0.4)) / 4
    for i, (v, lbl) in enumerate(stats):
        sx = x + Inches(0.2) + i * sw
        add_text(s, sx, y + Inches(0.55), sw, Inches(0.5),
                 v, size=24, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, sx, y + Inches(1.05), sw, Inches(0.3),
                 lbl, size=11, color=RGBColor(0xC8, 0xD0, 0xE0),
                 align=PP_ALIGN.CENTER)

    chart_y = y + Inches(1.6)
    add_picture(s, x, chart_y, w * 0.62, h - Inches(1.7),
                os.path.join(ASSETS, "chart_15_f1_baselines.png"))

    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    card(s, rx, chart_y, rw, h - Inches(1.7),
         icon="globe", title="On unseen code",
         body=("Held-out split:  F1 61.11 %\n"
               "External corpus:  F1 63.64 %\n"
               "→ recall gives ground, precision holds\n\n"
               "FAIRNESS\n"
               "FPR 10 %  ·  95 % CI [2.1–26.5 %]\n"
               "SAST shown at function scope,\nunderstating full-project recall (A7)\n\n"
               "Optional ≥ 0.67 rescore gate → A9"),
         body_size=11, title_size=14, accent_color=ACCENT, body_space=2)


# ============================================================
# SLIDE 16 — RAG ablation (with real chart)
# ============================================================
def slide_16_rag_ablation():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "RAG Is Not a Uniform Win", page_idx=16, show_brand=False)
    badge(s, SLIDE_W - Inches(2.4), Inches(0.25), Inches(1.8), Inches(0.4),
          "R1  ·  RQ2", color=ACCENT)
    x, y, w, h = content_box()
    add_picture(s, x, y, w * 0.62, h - Inches(0.2),
                os.path.join(ASSETS, "chart_16_rag_ablation.png"))
    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    add_rect(s, rx, y, rw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.2), y + Inches(0.15), rw - Inches(0.4), Inches(0.4),
             "TAKEAWAY", size=11, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.2), y + Inches(0.55), rw - Inches(0.4), Inches(2.5),
             "RAG can help, harm,\nor do nothing,\ndepending on the model.",
             size=20, bold=True, color=TEXT_LIGHT)
    add_text(s, rx + Inches(0.2), y + Inches(2.9), rw - Inches(0.4), Inches(1.8),
             "Only codellama's −29.89 F1 survives Bonferroni at α = 0.01\n(McNemar p = 0.0013).\n\nDesign implication:  treat RAG as a per-model configuration choice, not a default.",
             size=11, color=RGBColor(0xD8, 0xDE, 0xEE))


# ============================================================
# SLIDE 17 — R3 Repair (Huzefa-style bar metrics)
# ============================================================
def slide_17_r3_repair():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R3 · Repair Quality", page_idx=17, show_brand=False)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R3", color=ACCENT)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2
    panel_h = Inches(3.65)

    # Left column — Auto-applicable (objective / headline)
    card_bg(s, x, y, cw, panel_h, accent=ACCENT)
    place_icon(s, x + Inches(0.25), y + Inches(0.2), Inches(0.3), "gear")
    add_text(s, x + Inches(0.64), y + Inches(0.18), cw - Inches(0.9), Inches(0.4),
             "AUTO-APPLICABLE", size=13, bold=True, color=NAVY_DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.5), Inches(0.3),
             "Validator parses with @babel/parser  ·  fully automated",
             size=10, color=TEXT_MUTED)
    # Two big numbers as labeled bars, vertical with breathing room
    auto = [
        ("On issued fixes  (252 / 279)", 0.9032, "90.32 %", ACCENT),
        ("Across all stage-2 calls",     0.8485, "84.85 %", NAVY),
    ]
    my = y + Inches(1.05)
    for lbl, val, val_label, color in auto:
        add_text(s, x + Inches(0.25), my, cw - Inches(0.5), Inches(0.3),
                 lbl, size=11, bold=True, color=NAVY_DARK)
        bar_metric(s, x + Inches(0.25), my + Inches(0.32), cw - Inches(0.5), Inches(0.55),
                   label=val_label, value=val, color=color)
        my += Inches(1.15)

    # Right column — Manual rubric
    rx = x + cw + gap
    card_bg(s, rx, y, cw, panel_h, accent=NAVY)
    place_icon(s, rx + Inches(0.25), y + Inches(0.2), Inches(0.3), "user")
    add_text(s, rx + Inches(0.64), y + Inches(0.18), cw - Inches(0.9), Inches(0.4),
             "MANUAL REVIEW  ·  n = 25", size=13, bold=True, color=NAVY_DARK)
    add_text(s, rx + Inches(0.25), y + Inches(0.55), cw - Inches(0.5), Inches(0.3),
             "Rubric-based review on a sampled subset",
             size=10, color=TEXT_MUTED)
    rubric = [
        ("Fix provided",          0.76,  "76 %"),
        ("Semantically correct",  0.895, "89.5 %"),
        ("Strategy-aligned",      0.947, "94.7 %"),
        ("Executable code",       0.421, "42.1 %"),
    ]
    my = y + Inches(1.05)
    row_h = Inches(0.6)
    for lbl, val, val_label in rubric:
        add_text(s, rx + Inches(0.25), my, cw - Inches(0.5), Inches(0.22),
                 lbl, size=10, bold=True, color=NAVY_DARK)
        bar_metric(s, rx + Inches(0.25), my + Inches(0.24), cw - Inches(0.5), Inches(0.33),
                   label=val_label, value=val, color=NAVY)
        my += row_h

    by = y + panel_h + Inches(0.15)
    add_rect(s, x, by, w, Inches(0.95), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6), Inches(0.25),
             "TWO-TIER METRIC", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.36), w - Inches(0.6), Inches(0.5),
             "Auto-applicable = deployment reading.  Manual = decision-support reading.  Combined = 76 % fix-rate × 89.5 % semantic ≈ 68 %.  Stage 2 returns { code, language }; 18 of 297 calls returned no fix.",
             size=11, color=TEXT_LIGHT)


# ============================================================
# SLIDE 18 — R4 Latency (with real chart)
# ============================================================
def slide_18_r4_latency():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R4 · Usability: Latency Bands", page_idx=18, show_brand=False)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R4", color=ACCENT)
    x, y, w, h = content_box()
    add_picture(s, x, y, w * 0.62, h - Inches(0.2),
                os.path.join(ASSETS, "chart_18_latency_bands.png"))
    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    bands = [
        ("bolt", "Real-time  ≤ 500 ms", "SAST baselines only", RED),
        ("mouse", "Interactive  ≤ 1.5 s", "gemma3:1b + RAG  (1,019 ms)\nqwen3:4b  (1,328 ms)", ACCENT),
        ("chip", "On-demand  ≤ 5 s", "qwen3:8b + RAG  (2,216 ms)  ←  headline\np95 = 4,327 ms", GREEN),
    ]
    by = y
    bh = (h - Inches(0.4) - 2 * Inches(0.2)) / 3
    for iconname, hdr, body, col in bands:
        add_rect(s, rx, by, rw, bh, fill=CARD_BG, line=CARD_BORDER, rounded=True)
        add_rect(s, rx, by, Inches(0.1), bh, fill=col)
        place_icon(s, rx + Inches(0.22), by + Inches(0.12), Inches(0.3), iconname)
        add_text(s, rx + Inches(0.6), by + Inches(0.1), rw - Inches(0.7),
                 Inches(0.4), hdr, size=12, bold=True, color=NAVY_DARK)
        add_text(s, rx + Inches(0.2), by + Inches(0.45), rw - Inches(0.3),
                 bh - Inches(0.5), body, size=11, color=TEXT_DARK)
        by += bh + Inches(0.2)


# ============================================================
# SLIDE 19 — R5 Privacy (Huzefa-style bar metrics)
# ============================================================
def slide_19_r5_privacy():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R5 · Privacy: Empirically Verified", page_idx=19, show_brand=False)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R5", color=ACCENT)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2
    bottom_strip_h = Inches(0.55)
    avail = h - bottom_strip_h - Inches(0.15)
    ch = (avail - gap) / 2

    items = [
        ("lock", "Network Isolation",
         "Loopback-only by network policy. Source-leak\nscan: 0 non-loopback refs in 9,488 code windows.",
         "100 %  satisfied", 1.0),
        ("shield", "Prompt-Injection Harness",
         "12 crafted exfiltration attempts: no code\nechoed · schemaValidRate 83 % (10/12).",
         "100 %  leakFreeRate", 1.0),
        ("seal", "Ed25519-Signed Corpus",
         "RAG corpus provenance verified at\nload (manifest signature checked).",
         "Verified  ✓", 1.0),
        ("cube", "Reproducible Build",
         "Node 20.19.0-alpine · npm ci · seed = 42\nbyte-identical runs.",
         "Reproducible  ✓", 1.0),
    ]
    for i, (icon, t, b, val_label, val) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card_bg(s, cx, cy, cw, ch, accent=GREEN)
        # icon + title row
        place_icon(s, cx + Inches(0.22), cy + Inches(0.24), Inches(0.34), icon)
        add_text(s, cx + Inches(0.85), cy + Inches(0.22), cw - Inches(1.05), Inches(0.4),
                 t, size=14, bold=True, color=NAVY_DARK)
        # body block — anchored to space ABOVE the bar at the card bottom
        body_y = cy + Inches(0.72)
        bar_y = cy + ch - Inches(0.55)
        add_text(s, cx + Inches(0.22), body_y, cw - Inches(0.45),
                 bar_y - body_y - Inches(0.05), b, size=11, color=TEXT_DARK)
        # value bar pinned to bottom of card
        bar_metric(s, cx + Inches(0.22), bar_y, cw - Inches(0.45), Inches(0.42),
                   label=val_label, value=val, color=GREEN)

    by = y + 2 * ch + gap + Inches(0.05)
    add_rect(s, x, by, w, bottom_strip_h, fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by, w - Inches(0.6), bottom_strip_h,
             "Resource footprint:  harness 3 ms median CPU, under 100 MB RSS  ·  Ollama 6–8 GB RAM for 8B-q4.",
             size=11, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 20 — Takeaways
# ============================================================
def slide_20_takeaways():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "What the Evaluation Showed", page_idx=20)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(4.55) - gap) / 2
    items = [
        ("trophy", "A local 8B model is enough  ·  RQ1",
         "qwen3:8b + RAG reaches ~71 % F1 on JS/TS, well inside the on-demand latency band on a developer laptop. The privacy envelope holds without paying for it in quality."),
        ("book", "RAG is model-dependent  ·  RQ2",
         "The largest paired effect in the dataset is a RAG-induced degradation (codellama −29.89 F1). Treat retrieval as a per-model design choice, not a default."),
        ("wrench", "Repair is the real product gap",
         "90 % auto-applicable means the developer can act, not just be warned. This is where SAST falls flat: they detect, they don't fix."),
        ("lock", "Privacy as a verifiable boundary",
         "Loopback + signed corpus + container pin makes privacy testable, not aspirational. leakFreeRate = 100 % isn't a slogan; it's a measurement."),
    ]
    for i, (icon, t, b) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=icon, title=t, body=b,
             body_size=12, title_size=15, accent_color=NAVY,
             title_h=Inches(0.6), body_anchor=MSO_ANCHOR.MIDDLE)
    # explicit RQ closure — answer the three questions posed on slide 5
    by = y + 2 * ch + gap + Inches(0.12)
    add_rect(s, x, by, w, Inches(0.5), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.04), w - Inches(0.6), Inches(0.42),
             "RQ1  Yes: 71.43 % F1 under the privacy envelope     ·     "
             "RQ2  Yes, but per-model: RAG is not a default     ·     "
             "RQ3  Yes: 2.2 s median, inside the on-demand band",
             size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 21 — Code Guardian closes the gap (callback to slide 7)
# ============================================================
def slide_20b_comparison():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Code Guardian Closes the Gap", page_idx=21)
    x, y, w, h = content_box()

    rows = [
        ("GitHub Copilot",   ["✅","✅","✅","🟡","❌"]),
        ("Semgrep",          ["🟡","✅","❌","✅","✅"]),
        ("CodeQL",           ["🟡","✅","❌","🟡","✅"]),
        ("IRIS (GPT-4)",     ["🟡","🟡","❌","❌","❌"]),
        ("SecRepair",        ["🟡","🟡","🟡","❌","🟡"]),
        ("Code Guardian",    ["🟡","✅","✅","✅","✅"]),
    ]
    header = ["Solution / Approach", "R1\nACCURACY", "R2\nCONSISTENCY",
              "R3\nREPAIR", "R4\nUSABILITY", "R5\nPRIVACY"]

    table_w = w * 0.72
    legend_w = w - table_w - Inches(0.3)
    col_w = [Inches(2.4)] + [(table_w - Inches(2.4)) / 5] * 5

    # fill the content height: table + caption span the whole box
    caption_h = Inches(0.4)
    header_h = Inches(0.72)
    table_h = h - caption_h - Inches(0.12)
    row_h = (table_h - header_h) / len(rows)
    table_bottom = y + table_h

    cx = x
    for i, hd in enumerate(header):
        add_rect(s, cx, y, col_w[i], header_h, fill=NAVY)
        add_text(s, cx, y, col_w[i], header_h, hd,
                 size=10, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    ry = y + header_h
    for r_i, (name, cells) in enumerate(rows):
        is_self = name == "Code Guardian"
        bg = RGBColor(0xE9, 0xF3, 0xEA) if is_self else (LIGHT_BG if r_i % 2 == 0 else CARD_BG)
        cx = x
        add_rect(s, cx, ry, col_w[0], row_h, fill=bg, line=CARD_BORDER)
        if is_self:
            add_rect(s, cx, ry, Inches(0.1), row_h, fill=GREEN)
        add_text(s, cx + Inches(0.22), ry, col_w[0] - Inches(0.3), row_h,
                 name, size=12, bold=is_self, color=NAVY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[0]
        for c_i, cell in enumerate(cells):
            add_rect(s, cx, ry, col_w[c_i + 1], row_h, fill=bg, line=CARD_BORDER)
            cell_marker(s, cx, ry, col_w[c_i + 1], row_h, cell)
            cx += col_w[c_i + 1]
        ry += row_h

    add_text(s, x, table_bottom + Inches(0.06), table_w, caption_h,
             "Green check = meets the pre-registered threshold; amber = Medium.   "
             "R1 is amber: F1 71 % lands in the Medium band, competitive but not High.",
             size=9, italic=True, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # right column — legend + GAP box, bottom-aligned with the table
    lx = x + table_w + Inches(0.3)
    legend_h = Inches(1.7)
    add_rect(s, lx, y, legend_w, legend_h, fill=CARD_BG, line=CARD_BORDER,
             rounded=True)
    add_text(s, lx + Inches(0.2), y + Inches(0.14), legend_w - Inches(0.4),
             Inches(0.3), "LEGEND", size=10, bold=True, color=ACCENT)
    legend_rows(s, lx + Inches(0.2), y + Inches(0.55), legend_w - Inches(0.4),
                size=11)

    co_y = y + legend_h + Inches(0.22)
    co_h = table_bottom - co_y
    add_rect(s, lx, co_y, legend_w, co_h, fill=GREEN, rounded=True)
    add_text(s, lx + Inches(0.22), co_y + Inches(0.22), legend_w - Inches(0.44),
             Inches(0.35), "GAP CLOSED", size=12, bold=True, color=TEXT_LIGHT)
    add_text(s, lx + Inches(0.22), co_y + Inches(0.68), legend_w - Inches(0.44),
             co_h - Inches(0.9),
             "The only approach to bring privacy + repair + IDE-grade latency "
             "together in one local tool, with competitive detection accuracy.",
             size=13, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 22 — Limitations & Future Work (side by side)
# ============================================================
def slide_21_limitations_future():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Limitations & Future Work", page_idx=22)
    x, y, w, h = content_box()
    gap = Inches(0.4)
    col_w = (w - gap) / 2
    cols = [
        ("Current Limitations", RED, [
            ("user", "Single-reviewer manual review (n = 25)",
             "Headline R3 is the objective auto-applicable rate; inter-rater κ is future work."),
            ("target", "Curated 101-case corpus",
             "Offset by held-out 71 / 30 split, 15-case external set, whole-project NodeGoat run."),
            ("globe", "One real-world project end-to-end",
             "NodeGoat only  ·  multi-project validation pending."),
            ("refresh", "Consensus filter inert under determinism",
             "Three seed-42 runs are byte-identical; a real signal needs seed rotation."),
            ("cloud", "No cloud-LLM baseline",
             "Excluded by the privacy threat model  ·  3 SAST tools set the floor."),
        ]),
        ("Future Work", ACCENT, [
            ("dice", "Seed rotation",
             "Genuine consensus signal beyond the deterministic-decoding floor."),
            ("users", "Inter-rater agreement",
             "Second reviewer on the 25-sample subset for a κ statistic."),
            ("building", "Multi-project real-world validation",
             "Measured recall beyond NodeGoat, across diverse production codebases."),
            ("bolt", "Streaming partial results",
             "Surface first-token findings to break into the real-time band."),
            ("target", "Adaptive RAG gating",
             "Per-model RAG enable/disable driven by slide 16's finding."),
        ]),
    ]
    head_h = Inches(0.5)
    for c, (htitle, accent, items) in enumerate(cols):
        cx = x + c * (col_w + gap)
        add_text(s, cx, y, col_w, head_h, htitle, size=16, bold=True,
                 color=accent, anchor=MSO_ANCHOR.MIDDLE)
        rows_y = y + head_h + Inches(0.05)
        rh = (h - head_h - Inches(0.05)) / len(items)
        for i, (icon, t, b) in enumerate(items):
            ry = rows_y + i * rh
            add_rect(s, cx, ry, col_w, rh - Inches(0.1), fill=CARD_BG,
                     line=CARD_BORDER, rounded=True)
            add_rect(s, cx, ry, Inches(0.09), rh - Inches(0.1), fill=accent)
            isz = Inches(0.30)
            place_icon(s, cx + Inches(0.18), ry + (rh - Inches(0.1) - isz) / 2, isz,
                       icon, variant=("red" if accent == RED else "accent"))
            add_text(s, cx + Inches(0.72), ry + Inches(0.07), col_w - Inches(0.85),
                     Inches(0.33), t, size=12, bold=True, color=NAVY_DARK)
            add_text(s, cx + Inches(0.72), ry + Inches(0.38), col_w - Inches(0.85),
                     rh - Inches(0.45), b, size=10, color=TEXT_DARK)


# ============================================================
# SLIDE 23 — Contributions
# ============================================================
def slide_23_contributions():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Contributions of This Thesis", page_idx=23)
    x, y, w, h = content_box()
    items = [
        ("wrench", "Working VS Code extension",
         "Local LLM + RAG + signed corpus + pinned container, published to the Marketplace."),
        ("link", "Two-stage pipeline with structured contracts",
         "Detection (JSON-mode) + Repair ({ code, language }), schema-pinned at every boundary."),
        ("barchart", "Empirical study",
         "5 models × 2 modes + 3 SAST baselines · 303 invocations per config · CWE-mapped JS/TS corpus."),
        ("refresh", "Reproducibility infrastructure",
         "Byte-identical runs, Ed25519-signed manifest, pinned Node 20.19.0-alpine container."),
        ("ruler", "Statistical discipline",
         "Bonferroni-corrected McNemar, exact-binomial CIs, held-out 71 / 30 split."),
    ]
    rh = (h - Inches(0.2)) / len(items)
    for i, (icon, t, b) in enumerate(items):
        ry = y + i * rh
        add_rect(s, x, ry, w, rh - Inches(0.1), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_rect(s, x, ry, Inches(0.1), rh - Inches(0.1), fill=GREEN)
        isz = Inches(0.34)
        place_icon(s, x + Inches(0.2), ry + (rh - Inches(0.1) - isz) / 2, isz, icon)
        add_text(s, x + Inches(0.85), ry + Inches(0.08), Inches(4.5), Inches(0.4),
                 t, size=13, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.85), ry + Inches(0.45), w - Inches(1.0),
                 rh - Inches(0.5), b, size=11, color=TEXT_DARK)


# ============================================================
# SLIDE 24 — References (NEW)
# ============================================================
def slide_24_references():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "References", page_idx=24)
    x, y, w, h = content_box()
    refs = [
        "[1] Touvron, H., et al.: Llama 2: Open foundation and fine-tuned chat models. arXiv:2307.09288 (2023)",
        "[2] Lewis, P., et al.: Retrieval-augmented generation for knowledge-intensive NLP tasks. NeurIPS (2020)",
        "[3] Li, Z., et al.: IRIS: Using LLMs for context-aware vulnerability detection. ICLR (2025)",
        "[4] Islam, N., et al.: SecRepair: LLM-driven security vulnerability repair. NDSS (2024)",
        "[5] Kavian, S., et al.: LLMSecGuard: Securing LLM-generated code. EASE (2024)",
        "[6] Peng, B., et al.: CWEval: Outcome-driven evaluation of LLM code generation. ACL (2025)",
        "[7] Pearce, H., et al.: SecurityEval dataset: mining vulnerability examples. MSR4P&S (2022)",
        "[8] MITRE: Common Weakness Enumeration (CWE)  ·  https://cwe.mitre.org",
        "[9] OWASP Foundation: OWASP Top 10  ·  https://owasp.org/Top10",
        "[10] Malfa, E.L., Petrov, A., et al.: A taxonomy of evaluation methods for LLMs. arXiv:2402.13887 (2024)",
        "[11] Ollama: Local LLM serving  ·  https://ollama.com",
        "[12] Microsoft: VS Code Extension API  ·  https://code.visualstudio.com/api",
    ]
    # two columns
    col_w = (w - Inches(0.4)) / 2
    col_h = h - Inches(0.2)
    half = (len(refs) + 1) // 2
    for col in range(2):
        cx = x + col * (col_w + Inches(0.4))
        ry = y
        items = refs[col * half: (col + 1) * half]
        for r in items:
            add_text(s, cx, ry, col_w, Inches(0.6), r, size=10, color=TEXT_DARK,
                     italic=True)
            ry += Inches(0.55)


# ============================================================
# SLIDE 25 — Thank You
# ============================================================
def slide_25_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(0.6), Inches(0.7), Inches(0.04), Inches(0.5), fill=ACCENT)
    add_text(s, Inches(0.85), Inches(0.65), Inches(8), Inches(0.5),
             "TECHNICAL UNIVERSITY OF CHEMNITZ", size=12, bold=True,
             color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(2.2), Inches(6.5), Inches(1.4),
             "Thank You", size=72, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(3.7), Inches(6.5), Inches(0.8),
             "Questions & Discussion", size=30,
             color=RGBColor(0xC8, 0xD0, 0xE0))

    # right-side decorative question bubble
    qbubble_x = Inches(8.5)
    add_rect(s, qbubble_x, Inches(2.4), Inches(3.8), Inches(2.6),
             fill=NAVY_DARK, rounded=True)
    add_text(s, qbubble_x, Inches(2.4), Inches(3.8), Inches(2.6),
             "?", size=140, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    strip_y = Inches(5.4)
    add_rect(s, Inches(0.85), strip_y, SLIDE_W - Inches(1.7), Inches(0.85),
             fill=NAVY_DARK, rounded=True)
    add_text(s, Inches(1.0), strip_y + Inches(0.1), SLIDE_W - Inches(2.0),
             Inches(0.3), "HEADLINE", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), strip_y + Inches(0.35), SLIDE_W - Inches(2.0),
             Inches(0.4),
             "F1 71.43 %   ·   Repair 90.32 % auto-applicable   ·   Latency 2.2 s median   ·   leakFreeRate 100 %",
             size=13, bold=True, color=TEXT_LIGHT)

    by = Inches(6.55)
    add_text(s, Inches(0.85), by, Inches(4.5), Inches(0.3),
             "Md Hafizur Rahman", size=12, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), by + Inches(0.28), Inches(4.5), Inches(0.3),
             "M.Sc. Web Engineering  ·  TU Chemnitz",
             size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
    add_text(s, Inches(9.5), by, Inches(3.5), Inches(0.3),
             "DEFENSE  ·  " + DATE_STR, size=10, bold=True, color=ACCENT,
             align=PP_ALIGN.RIGHT)


# ============================================================
# APPENDIX
# ============================================================
def appendix_header(slide, title):
    add_slide_chrome(slide, title, page_idx=None, show_brand=False)
    badge(slide, SLIDE_W - Inches(1.8), Inches(0.25), Inches(1.2), Inches(0.4),
          "APPENDIX", color=ACCENT)


# Canonical per-model results filled from
# evaluation/logs/evaluation-2026-04-25T11-44-45-606Z-ablation-baselines.json
# (canonical taxonomy) and phase1-b3 for qwen3:8b+RAG headline.
def slide_A1_full_results():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A1  ·  Full Per-Model Results  ·  Canonical Taxonomy")
    x, y, w, h = content_box()
    headers = ["Configuration", "F1", "Prec", "Recall", "FPR",
               "Latency p50", "VRAM"]
    rows = [
        ("qwen3:8b + RAG",       "71.43", "72.46", "70.42", "10.00", "2,216 ms", "~6 GB"),
        ("qwen3:8b (no RAG)",    "67.52", "61.63", "74.65", "50.00", "1,534 ms", "~6 GB"),
        ("qwen3:4b + RAG",       "60.92", "51.46", "74.65", "100.00", "1,328 ms", "~4 GB"),
        ("qwen3:4b (no RAG)",    "66.14", "56.95", "78.87", "90.00", "1,305 ms", "~4 GB"),
        ("gemma3:4b + RAG",      "57.14", "46.85", "73.24", "100.00", "2,077 ms", "~4 GB"),
        ("gemma3:4b (no RAG)",   "57.30", "46.49", "74.65", "100.00", "2,550 ms", "~4 GB"),
        ("gemma3:1b + RAG",      "36.02", "32.22", "40.85", "23.33", "1,019 ms", "~2 GB"),
        ("gemma3:1b (no RAG)",   "36.65", "29.17", "49.30", "53.33", "1,215 ms", "~2 GB"),
        ("codellama + RAG",      "23.17", "14.79", "53.52", "100.00", "3,529 ms", "~6 GB"),
        ("codellama (no RAG)",   "53.06", "41.60", "73.24", "100.00", "2,024 ms", "~6 GB"),
        ("Semgrep (SAST)",       "20.22", "50.00", "12.68", "6.67",  "63 ms",   "n/a"),
        ("CodeQL (SAST)",        "10.85", "12.07", "9.86",  "53.33", "182 ms",  "n/a"),
        ("ESLint-security",      "0.00",  "0.00",  "0.00",  "0.00",  "1 ms",    "n/a"),
    ]
    col_w = [Inches(2.3), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95),
             Inches(1.5), Inches(0.95)]
    row_h = Inches(0.32)
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, col_w[i], Inches(0.4), fill=NAVY)
        add_text(s, cx, y, col_w[i], Inches(0.4), hd, size=10, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]
    ry = y + Inches(0.4)
    for r_i, row in enumerate(rows):
        is_head = row[0] == "qwen3:8b + RAG"
        bg = RGBColor(0xE9, 0xF3, 0xEA) if is_head else (LIGHT_BG if r_i % 2 == 0 else CARD_BG)
        cx = x
        for c_i, val in enumerate(row):
            add_rect(s, cx, ry, col_w[c_i], row_h, fill=bg, line=CARD_BORDER)
            add_text(s, cx, ry, col_w[c_i], row_h, val, size=9.5,
                     bold=(c_i == 0 or is_head),
                     color=NAVY_DARK,
                     align=PP_ALIGN.CENTER if c_i > 0 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c_i]
        ry += row_h
    add_text(s, x, ry + Inches(0.05), w, Inches(0.3),
             "qwen3:8b + RAG headline uses phase1-b3 canonical run; all others use canonical-taxonomy ablation. p95 latencies in source JSON.",
             size=9, color=TEXT_MUTED, italic=True)


# Confusion matrices derived from canonical-taxonomy qwen3:8b+RAG run.
# 71 vulnerable cases, 30 secure cases. Run-level results expanded back
# to case-level using ≥2-of-3 majority detection.
def slide_A2_confusion():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A2  ·  Confusion Matrices  ·  qwen3:8b + RAG")
    x, y, w, h = content_box()
    gap = Inches(0.4)
    cw = (w - gap) / 2

    def matrix(cx, cy, label, tp, fn, fp, tn):
        add_rect(s, cx, cy, cw, Inches(3.8), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.25), cy + Inches(0.15), cw - Inches(0.5),
                 Inches(0.4), label, size=14, bold=True, color=NAVY)
        gx = cx + Inches(0.6)
        gy = cy + Inches(0.95)
        cw2 = (cw - Inches(1.2)) / 2
        rh = Inches(1.0)
        add_text(s, gx, cy + Inches(0.6), cw2, Inches(0.3),
                 "Predicted: Vuln", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
        add_text(s, gx + cw2, cy + Inches(0.6), cw2, Inches(0.3),
                 "Predicted: Safe", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
        add_text(s, cx + Inches(0.05), gy, Inches(0.55), rh,
                 "Actual\nVuln", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.05), gy + rh, Inches(0.55), rh,
                 "Actual\nSafe", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cells = [(tp, GREEN, "TP"), (fn, RED, "FN"),
                 (fp, RED, "FP"), (tn, GREEN, "TN")]
        for i, (v, col, lbl) in enumerate(cells):
            r, c = i // 2, i % 2
            ccx = gx + c * cw2
            ccy = gy + r * rh
            bg = RGBColor(0xE9, 0xF3, 0xEA) if col == GREEN else RGBColor(0xF9, 0xE6, 0xEA)
            add_rect(s, ccx, ccy, cw2, rh, fill=bg, line=CARD_BORDER)
            add_text(s, ccx, ccy + Inches(0.05), cw2, Inches(0.4),
                     str(v), size=22, bold=True, color=NAVY_DARK,
                     align=PP_ALIGN.CENTER)
            add_text(s, ccx, ccy + Inches(0.55), cw2, Inches(0.3),
                     lbl, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # case-level (n=101): TP = vulnerable cases majority-detected = ~50, FN = 21,
    # FP = secure cases with at least one FP = 3 (gated), TN = 27.
    # These are derived from secureFalsePositiveCases (gated 9 runs → 3 cases ≥2/3)
    # and the 71 vulnerable cases majority-detection.
    matrix(x, y, "Full corpus  (n = 101)  ·  case-level majority", "50", "21", "3", "27")
    # Held-out TEST split = 21 vulnerable + 9 secure (Section 5.4.7): recall 52.38 % → TP=11,
    # FN=10; FPR 0 % → FP=0, TN=9.
    matrix(x + cw + gap, y, "Held-out test split  (n = 30)", "11", "10", "0", "9")
    add_text(s, x, y + Inches(3.95), w, Inches(0.3),
             "Sample-level via ≥2/3 majority over 3 deterministic runs (seed=42, T=0).  Precision is finding-level/pooled, so 72.46 % ≠ TP/(TP+FP) here.  Source: phase1-b3 canonical taxonomy.",
             size=9, color=TEXT_MUTED, italic=True)


def slide_A3_manual_review():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A3  ·  Why the Single-Reviewer Manual Review Is Bounded")
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(3.6) - gap) / 2
    rubric = [
        ("Fix provided", "76 %", "Of 25 reviewed samples, 19 included a repair suggestion."),
        ("Semantically correct", "89.5 %", "Fix addresses the CWE (17 of 19 issued fixes)."),
        ("Strategy-aligned", "94.7 %", "Right repair pattern (parameterized query, escaping, etc.)."),
        ("Executable code", "42.1 %", "Patch runs without modification; the rest give prose guidance."),
    ]
    for i, (lbl, val, sub) in enumerate(rubric):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card_bg(s, cx, cy, cw, ch, accent=NAVY)
        add_text(s, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4),
                 Inches(0.4), lbl, size=13, bold=True, color=NAVY_DARK)
        add_text(s, cx + Inches(0.2), cy + Inches(0.55), cw - Inches(0.4),
                 Inches(0.7), val, size=30, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.2), cy + Inches(1.3), cw - Inches(0.4),
                 ch - Inches(1.5), sub, size=10, color=TEXT_DARK)
    by = y + 2 * ch + gap + Inches(0.1)
    add_rect(s, x, by, w, Inches(0.7), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6), Inches(0.5),
             "Inter-rater on the 25-sample subset = future work.  Headline R3 is the fully objective auto-applicable rate (90.32 %, n = 279).",
             size=12, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_A4_corpus():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A4  ·  Why a Curated Corpus, Not an Existing Benchmark")
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2

    add_rect(s, x, y, cw, h - Inches(0.2), fill=CARD_BG,
             line=CARD_BORDER, rounded=True)
    add_text(s, x + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4), "WHY EXISTING JS/TS BENCHMARKS DON'T FIT", size=11,
             bold=True, color=ACCENT)
    items = [
        ("Juliet (NIST)", "C / C++ / Java only.  No first-party JS port."),
        ("OWASP Benchmark", "Java only."),
        ("SecurityEval (Pearce 2022)", "Small, partially synthetic.  CWE labelling inconsistent with OWASP Top 10."),
        ("CyberSecEval", "29.3 % of vulnerable samples are reproducible (Peng et al. 2025)."),
        ("CWEval (Peng 2025)", "119 curated tasks across 31 CWEs, 5 languages, peer-comparable to 101 cases."),
    ]
    ry = y + Inches(0.6)
    for name, sub in items:
        add_text(s, x + Inches(0.25), ry, cw - Inches(0.5), Inches(0.3),
                 name, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), ry + Inches(0.3), cw - Inches(0.5),
                 Inches(0.45), sub, size=10, color=TEXT_DARK)
        ry += Inches(0.75)

    rx = x + cw + gap
    add_rect(s, rx, y, cw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4),
             "SELECTION-BIAS MITIGATIONS IN THIS THESIS", size=11, bold=True,
             color=ACCENT)
    mits = [
        ("Held-out 71 / 30 split", "Frozen TEST set keeps threshold tuning out of headline numbers."),
        ("15-case external corpus", "NodeGoat, Juice Shop, 3 named CVEs, independent of the curator."),
        ("Whole-project NodeGoat run", "Project-level recall independent of sample selection."),
        ("Provenance bound into manifest", "Ed25519 signature; corpus cannot be silently rewritten."),
    ]
    ry = y + Inches(0.65)
    for name, sub in mits:
        add_text(s, rx + Inches(0.25), ry, cw - Inches(0.5), Inches(0.3),
                 name, size=12, bold=True, color=TEXT_LIGHT)
        add_text(s, rx + Inches(0.25), ry + Inches(0.3), cw - Inches(0.5),
                 Inches(0.55), sub, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
        ry += Inches(0.95)


# Threat-model now uses real diagram
# ---- native (editable) threat-model diagram -------------------------------
# Ported 1:1 from make_assets.diagram_threat_model (100x100 user-unit canvas).
BOUND_GREEN_FILL = RGBColor(0xEC, 0xF6, 0xF0)
TAN_FILL = RGBColor(0xF3, 0xE9, 0xD8)
CLOUD_FILL = RGBColor(0xFD, 0xE9, 0xEE)
SUB_ON_NAVY2 = RGBColor(0xC2, 0xCB, 0xDD)


def draw_threat_model(slide, x, y, w, h):
    # map matplotlib content region -> the given EMU box (crop bare margins)
    MXL, MXR, MYB, MYT = 0.0, 100.0, 5.0, 96.0
    SCALE = (w / 914400.0) / 6.6

    def MX(v): return int(round(x + (v - MXL) / (MXR - MXL) * w))
    def MY(v): return int(round(y + (1 - (v - MYB) / (MYT - MYB)) * h))
    def MW(v): return int(round(v / (MXR - MXL) * w))
    def MH(v): return int(round(v / (MYT - MYB) * h))
    def FS(base): return Pt(base * SCALE)

    def text_at(cx, cy, lines, color=TEXT_DARK, size=9.0, bold=False,
                italic=False, align=PP_ALIGN.CENTER, boxw=40, anchor=MSO_ANCHOR.MIDDLE):
        if isinstance(lines, str):
            lines = [lines]
        cw = MW(boxw); ch = Pt(size * SCALE * 1.5 * len(lines) + 4)
        if align == PP_ALIGN.CENTER:
            left = MX(cx) - cw // 2
        elif align == PP_ALIGN.RIGHT:
            left = MX(cx) - cw
        else:
            left = MX(cx)
        tb = slide.shapes.add_textbox(left, MY(cy) - int(ch // 2), cw, int(ch))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = anchor
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = ln
            r.font.size = FS(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; r.font.name = FONT_BODY
        tb.fill.background(); tb.line.fill.background()
        tb.shadow.inherit = False
        return tb

    def badge(cx, cy, glyph, col):
        d = MH(5.6)
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, MX(cx) - d // 2,
                                    MY(cy) - d // 2, d, d)
        set_solid(ov, col); no_line(ov); ov.shadow.inherit = False
        tf = ov.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = glyph
        r.font.size = FS(10.5); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT; r.font.name = FONT_BODY

    def card(bx, by, bw, bh, title, sub, fill=CARD_BG, edge=NAVY, tcol=NAVY_DARK,
             scol=TEXT_MUTED, mark=None, mark_col=None, lw=1.6, ts=11.5):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX(bx),
                                    MY(by + bh), MW(bw), MH(bh))
        try:
            sp.adjustments[0] = 0.12
        except Exception:
            pass
        set_solid(sp, fill); set_line(sp, edge, lw); sp.shadow.inherit = False
        # text — shifted right of the badge when present
        tlines = [l for l in (sub.split("\n") if sub else []) if l != ""]
        lx = bx + (9 if mark else 1)
        rx = bx + bw - 1
        cx = (lx + rx) / 2
        tf_left = MX(lx); tf_w = MX(rx) - MX(lx)
        tb = slide.shapes.add_textbox(tf_left, MY(by + bh), tf_w, MH(bh))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = FS(ts); r.font.bold = True
        r.font.color.rgb = tcol; r.font.name = FONT_BODY
        for ln in tlines:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(1)
            r2 = p2.add_run(); r2.text = ln
            r2.font.size = FS(8.5); r2.font.color.rgb = scol
            r2.font.name = FONT_BODY
        tb.fill.background(); tb.line.fill.background(); tb.shadow.inherit = False
        if mark:
            badge(bx + 6, by + bh / 2, mark, mark_col)

    # ---- privacy boundary (back) ----
    bnd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MX(2), MY(94),
                                 MW(66), MH(77))
    try:
        bnd.adjustments[0] = 0.04
    except Exception:
        pass
    set_solid(bnd, BOUND_GREEN_FILL); set_line(bnd, GREEN, 2.0)
    line_dash(bnd.line); bnd.shadow.inherit = False
    text_at(35, 90, "YOUR LAPTOP  ·  privacy boundary", color=GREEN, size=10.0,
            bold=True, boxw=62)

    # ---- inside: stays local (green check) ----
    card(7, 70, 52, 13, "Code Guardian", "local LLM + RAG", fill=NAVY, edge=NAVY,
         tcol=TEXT_LIGHT, scol=SUB_ON_NAVY2, mark="✓", mark_col=GREEN)
    card(7, 53, 52, 13, "Local SAST baselines", "Semgrep · CodeQL · ESLint",
         mark="✓", mark_col=GREEN)
    card(7, 36, 52, 13, "Your code & prompts", "source · embeddings · context",
         mark="✓", mark_col=GREEN)
    card(7, 23, 52, 8, "Only outbound: signed corpus updates", "",
         fill=TAN_FILL, edge=TEXT_MUTED, tcol=TEXT_MUTED, ts=9.5, lw=1.2)

    # ---- outside: excluded cloud baseline (red cross) ----
    card(73, 36, 25, 13, "Cloud LLM", "Copilot · GPT-4\nClaude · Cursor",
         fill=CLOUD_FILL, edge=RED, tcol=RED, scol=RED, ts=11, mark="✗",
         mark_col=RED)

    # ---- egress a cloud baseline would require ----
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, MX(59), MY(42.5),
                                   MX(73), MY(42.5))
    c.line.color.rgb = RED; c.line.width = Pt(1.8)
    line_dash(c.line); line_arrow(c.line); c.shadow.inherit = False
    text_at(66, 48, ["requires", "code egress"], color=RED, size=8.0, bold=True,
            italic=True, boxw=20)

    text_at(50, 9.5, "Cloud baselines can't be run without breaking the "
            "threat model.", color=RED, size=8.5, bold=True, boxw=96)


def slide_A5_no_cloud():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A5  ·  Why No Cloud-LLM Baseline")
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2

    draw_threat_model(s, x, y, cw, h - Inches(0.2))

    rx = x + cw + gap
    add_rect(s, rx, y, cw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4), "ARGUMENT", size=11, bold=True, color=ACCENT)
    pts = [
        ("Cloud-LLM is the problem statement",
         "Including it as a baseline contradicts the threat model."),
        ("Floor provided by 3 SAST baselines",
         "Semgrep, CodeQL, ESLint  ·  all local, all reproducible."),
        ("Frontier cloud will always win raw F1",
         "The thesis competes on the constraint set:  local, zero-egress, reproducible."),
        ("Post-defense addendum",
         "A one-paragraph comparison to published GPT-4 numbers on similar JS corpora is straightforward."),
    ]
    ry = y + Inches(0.6)
    for name, sub in pts:
        add_text(s, rx + Inches(0.25), ry, cw - Inches(0.5), Inches(0.35),
                 name, size=12, bold=True, color=TEXT_LIGHT)
        add_text(s, rx + Inches(0.25), ry + Inches(0.35), cw - Inches(0.5),
                 Inches(0.7), sub, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
        ry += Inches(1.05)


def slide_A6_consensus():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A6  ·  Why R2 Reports 100 % Inter-Run Agreement")
    x, y, w, h = content_box()
    gap = Inches(0.2)
    cw = (w - 2 * gap) / 3

    json_text = ('{\n  "findings": [\n    {\n      "category":\n        "sql-injection",\n'
                 '      "severity":\n        "high",\n      "line": 42\n    }\n  ]\n}')
    for i in range(3):
        cx = x + i * (cw + gap)
        add_rect(s, cx, y, cw, Inches(3.5),
                 fill=NAVY_DARK, line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.2), y + Inches(0.1), cw - Inches(0.4),
                 Inches(0.35), f"RUN {i + 1}", size=10, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.2), y + Inches(0.45), cw - Inches(0.4),
                 Inches(3.0), json_text, size=11, color=TEXT_LIGHT,
                 font="Courier New")

    by = y + Inches(3.7)
    add_rect(s, x, by, w, Inches(0.95), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6),
             Inches(0.3),
             "WHY THEY'RE IDENTICAL", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.4), w - Inches(0.6),
             Inches(0.5),
             "temperature = 0   ·   seed = 42   ·   format = JSON   ·   fixed prompt.  Byte-identical → consensus signal requires seed rotation (future work).",
             size=12, color=TEXT_LIGHT)


# ============================================================
# SLIDE A7 — SAST baseline configuration (defends the comparison)
# ============================================================
def slide_A7_sast_config():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A7  ·  SAST Baseline Configuration")
    x, y, w, h = content_box()

    rows = [
        ("Semgrep", "20.22 %",
         "Registry packs:  p/security-audit · p/javascript ·\np/typescript · p/owasp-top-ten   (--metrics=off)"),
        ("CodeQL", "10.85 %",
         "Suite:  javascript-security-and-quality.qls\nJavaScript database, default queries"),
        ("ESLint", "0.00 %",
         "eslint-plugin-security v4  ·  recommended config\n(eslint ^9.21); no rule fired on the corpus"),
    ]
    rh = Inches(0.95)
    for i, (tool, f1, cfg) in enumerate(rows):
        ry = y + i * (rh + Inches(0.12))
        add_rect(s, x, ry, w, rh, fill=CARD_BG, line=CARD_BORDER, rounded=True)
        add_rect(s, x, ry, Inches(0.1), rh, fill=TEXT_MUTED)
        add_text(s, x + Inches(0.25), ry, Inches(2.0), rh,
                 tool, size=15, bold=True, color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(2.2), ry, Inches(1.4), rh,
                 f1, size=15, bold=True, color=RED, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(3.8), ry + Inches(0.16), w - Inches(4.1), rh - Inches(0.3),
                 cfg, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    by = y + 3 * (rh + Inches(0.12)) + Inches(0.1)
    add_rect(s, x, by, w, Inches(1.15), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.12), w - Inches(0.6), Inches(0.3),
             "FAIRNESS CAVEAT", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.45), w - Inches(0.6), Inches(0.6),
             "All three were run per-snippet, at the same function-level scope as the LLM, the only "
             "scope-parity setting. In full-project deployment they get cross-file data flow, so these "
             "recall figures understate their production performance. Read as “SAST under LLM-equivalent scope.”",
             size=12, color=TEXT_LIGHT)


# ============================================================
# SLIDE A8 — Requirement scales: metrics, thresholds, sources
# (defends "how did you choose the level thresholds?")
# ============================================================
def slide_A8_requirement_scales():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A8  ·  Requirement Scales: Metrics, Thresholds & Sources")
    x, y, w, h = content_box()

    headers = ["Requirement", "Metric(s)", "Levels & thresholds", "Source"]
    col_w = [w * 0.16, w * 0.22, w * 0.42, w * 0.20]
    rows = [
        ("R1 · Accuracy", "F1, precision, recall, FPR",
         "High: F1≥.75, P/R≥.70, FPR≤.20  ·  Medium: F1≥.60, FPR≤.40  ·  "
         "Low: F1≥.45  ·  Insufficient: F1<.45 or FPR>.70",
         "Johnson 2013;\nJi 2023"),
        ("R2 · Consistency", "JSON parse rate, inter-run agreement",
         "High: parse≥99%, agreement≥90%  ·  Medium: 95–99%  ·  "
         "Low: 85–95%  ·  None: <85%",
         "Harness\n(raw agreement)"),
        ("R3 · Repair", "Auto-applicable parse rate + manual rubric",
         "High ≥75%  ·  Medium 50–75%  ·  Low <50%   "
         "(3 bands; repair quality is partly qualitative)",
         "Monperrus 2014"),
        ("R4 · Usability", "End-to-end median latency",
         "High: ≤5 s inline / ≤15 s full-file  ·  Medium: 5–15 s  ·  Low: >15 s",
         "Card 1983;\nNielsen 1994"),
        ("R5 · Privacy", "5-item verification checklist",
         "Binary: all five criteria pass, or the requirement fails "
         "(privacy is a design constraint, not a spectrum)",
         "Design constraint\n(Ch. 2)"),
    ]

    header_h = Inches(0.4)
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, col_w[i], header_h, fill=NAVY)
        add_text(s, cx, y, col_w[i], header_h, hd, size=11, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    row_h = Inches(0.72)
    ry = y + header_h
    for r_i, row in enumerate(rows):
        bg = LIGHT_BG if r_i % 2 == 0 else CARD_BG
        cx = x
        for c_i, val in enumerate(row):
            add_rect(s, cx, ry, col_w[c_i], row_h, fill=bg, line=CARD_BORDER)
            add_text(s, cx + Inches(0.04), ry, col_w[c_i] - Inches(0.08), row_h,
                     val, size=(10.5 if c_i == 0 else 9.5),
                     bold=(c_i == 0), color=NAVY_DARK,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c_i]
        ry += row_h

    by = ry + Inches(0.12)
    add_rect(s, x, by, w, Inches(0.85), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6), Inches(0.3),
             "WHY THESE THRESHOLDS", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.38), w - Inches(0.6), Inches(0.42),
             "Fixed in Chapter 2 before any model ran, the same pre-registration discipline as the held-out split.  "
             "Accuracy and repair bands follow IR/SE convention; latency bands follow HCI response-time limits; "
             "privacy is a binary design constraint.",
             size=11, color=TEXT_LIGHT)


# ============================================================
# SLIDE A9 — Confidence gate (eval-only rescore filter)
# (defends the 74.07 % number moved off slide 15)
# ============================================================
def slide_A9_confidence_gate():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A9  ·  Confidence Gate: Eval-Only Rescore Filter")
    x, y, w, h = content_box()

    # before / after table
    headers = ["Metric", "Ungated", "Gate ≥ 0.67", ""]
    col_w = [w * 0.28, w * 0.22, w * 0.22, w * 0.28]
    rows = [
        ("F1",        "71.43 %", "74.07 %", "▲ precision-driven"),
        ("Precision", "72.46 %", "78.13 %", "▲ +5.67 pts"),
        ("Recall",    "70.42 %", "70.42 %", "= unchanged"),
        ("FPR",       "10.0 %",  "10.0 %",  "= unchanged"),
    ]
    header_h = Inches(0.4)
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, col_w[i], header_h, fill=NAVY)
        add_text(s, cx, y, col_w[i], header_h, hd, size=11, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]
    row_h = Inches(0.5)
    ry = y + header_h
    for r_i, row in enumerate(rows):
        bg = LIGHT_BG if r_i % 2 == 0 else CARD_BG
        cx = x
        for c_i, val in enumerate(row):
            add_rect(s, cx, ry, col_w[c_i], row_h, fill=bg, line=CARD_BORDER)
            add_text(s, cx + Inches(0.04), ry, col_w[c_i] - Inches(0.08), row_h,
                     val, size=(11 if c_i == 0 else 11),
                     bold=(c_i == 0), color=NAVY_DARK,
                     align=(PP_ALIGN.LEFT if c_i in (0, 3) else PP_ALIGN.CENTER),
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c_i]
        ry += row_h

    by = ry + Inches(0.2)
    add_rect(s, x, by, w, Inches(2.0), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.12), w - Inches(0.6), Inches(0.3),
             "WHAT THE GATE ACTUALLY DOES", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.45), w - Inches(0.6), Inches(1.5),
             "• Drops findings the model could not map to a specific canonical category (the low-signal "
             "'other' bucket): 6 findings in this run, all false positives. Recall is untouched.\n"
             "• NOT run consensus:  the three evaluation runs share seed 42 and are byte-identical, so there is "
             "no run-to-run disagreement to gate on, so the gate reduces to a category filter.\n"
             "• Off by default in the shipped extension (minConfidence).  Both gated and ungated numbers are "
             "reported in the thesis, so the reader picks the bar.",
             size=11, color=TEXT_LIGHT, space_after=4)


# ============================================================
# BUILD
# ============================================================
slide_01_title()
slide_02_problem_illustration()
slide_03_two_failures()
slide_04_three_gap()
slide_05_rqs_reqs()
slide_06_existing_approaches()
slide_07_gap_matrix()
slide_08_workflows()
slide_09_two_stage()
slide_10_screenshot()
slide_11_architecture()
slide_12_privacy_cards()
slide_13_eval_setup()
slide_14_hero()
slide_15_r1_accuracy()
slide_16_rag_ablation()
slide_17_r3_repair()
slide_18_r4_latency()
slide_19_r5_privacy()
slide_20_takeaways()
slide_20b_comparison()
slide_21_limitations_future()
slide_23_contributions()
slide_24_references()
slide_25_thanks()
slide_A1_full_results()
slide_A2_confusion()
slide_A3_manual_review()
slide_A4_corpus()
slide_A5_no_cloud()
slide_A6_consensus()
slide_A7_sast_config()
slide_A8_requirement_scales()
slide_A9_confidence_gate()

def _add_md_runs(p, text, size=11, color=TEXT_DARK):
    """Render a markdown line into a paragraph: **bold** spans become bold runs;
    stray single-* italic markers are stripped."""
    for i, seg in enumerate(text.split("**")):
        bold = i % 2 == 1
        if not bold:
            seg = seg.replace("*", "")
        if seg == "":
            continue
        r = p.add_run(); r.text = seg
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = FONT_BODY


def attach_speaker_notes(presentation, md_path):
    """Parse SPEAKER_NOTES.md and write each '## Slide N — …' block into the
    notes pane of slide N (1-indexed → presentation.slides[N-1])."""
    if not os.path.exists(md_path):
        print(f"!! speaker notes not found: {md_path}")
        return
    with open(md_path, encoding="utf-8") as f:
        lines = f.read().split("\n")

    blocks, cur, header, buf = {}, None, None, []

    def flush():
        if cur is not None:
            blocks[cur] = (header, list(buf))

    for ln in lines:
        m = re.match(r"^## Slide (\d+)\b", ln)
        if m:
            flush()
            cur = int(m.group(1)); header = ln[3:].strip(); buf = []
            continue
        if ln.startswith("## "):           # a non-slide section ends the notes
            flush(); cur = None; continue
        if cur is not None and ln.strip() != "---":
            buf.append(ln)
    flush()

    attached = 0
    for num, (hdr, body) in blocks.items():
        idx = num - 1
        if not (0 <= idx < len(presentation.slides)):
            continue
        tf = presentation.slides[idx].notes_slide.notes_text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = hdr
        r0.font.bold = True; r0.font.size = Pt(12); r0.font.name = FONT_BODY
        for raw in body:
            if raw.strip() == "":
                continue
            p = tf.add_paragraph(); p.space_after = Pt(4)
            _add_md_runs(p, raw.strip())
        attached += 1
    print(f"OK  attached speaker notes to {attached} slides")


attach_speaker_notes(prs, os.path.join(os.path.dirname(__file__),
                                       "SPEAKER_NOTES.md"))

total = 25
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if "/ __" in r.text:
                        r.text = r.text.replace("/ __", f"/ {total}")

out = os.path.join(os.path.dirname(__file__), "defense_slides_v3.pptx")
prs.save(out)
print(f"OK  wrote {out}  ·  {len(prs.slides)} slides")
