"""
Builds defense_slides_v2.pptx from the REDESIGN_PLAN.md spec.

Design language mirrors Huzefa's INVOX deck:
- Dark navy header band, white/light body
- Card-based content (rounded rects with icon + title + body)
- One idea per slide, hero-number cards on results slides
- Placeholder rectangles where the user must drop in screenshots / diagrams / charts

After running, open in PowerPoint/Keynote, drop in:
  Slide 2:  illustration  (problem framing image)
  Slide 10: VS Code extension screenshot
  Slide 11: architecture diagram (use thesis privacy_boundary.drawio.pdf as start)
  Slide 15: F1 bar chart (Code Guardian vs SAST baselines)
  Slide 16: RAG ablation bar chart (5 models, +/- delta)
  Slide 18: latency band bar chart
  A1: full per-model results table
  A2: confusion matrices
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ---- design tokens ----
NAVY = RGBColor(0x14, 0x21, 0x3D)
NAVY_DARK = RGBColor(0x0B, 0x14, 0x28)
ACCENT = RGBColor(0xFC, 0xA3, 0x11)
RED = RGBColor(0xC9, 0x18, 0x4A)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xFB)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE6, 0xEE)
TEXT_DARK = RGBColor(0x1A, 0x1F, 0x36)
TEXT_MUTED = RGBColor(0x5B, 0x65, 0x7A)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9 widescreen
HEADER_H = Inches(0.9)
FOOTER_H = Inches(0.35)
MARGIN = Inches(0.5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

DATE_STR = "11.05.2026"
AUTHOR_FOOTER = f"Md Hafizur Rahman  ·  Code Guardian  ·  TU Chemnitz  ·  {DATE_STR}"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

slides_meta = []  # collected for page numbering after build


def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=CARD_BG, line=None, line_pt=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        # tame the round-corner radius
        try:
            s.adjustments[0] = 0.06
        except Exception:
            pass
    set_solid(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, line_pt)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs: list of (text, {size, bold, color, italic}) tuples; '\\n' starts new para."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    first = True
    for text, opts in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            first = True
            continue
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", FONT_BODY)
        r.font.size = Pt(opts.get("size", 14))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", TEXT_DARK)
    return tb


def add_slide_chrome(slide, title, *, page_idx=None, header=True, footer=True,
                     show_title=True, accent_strip=True):
    # full white background
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
    if header:
        hdr = add_rect(slide, 0, 0, SLIDE_W, HEADER_H, fill=NAVY)
        # TU CHEMNITZ left
        add_text(slide, Inches(0.4), Inches(0.22), Inches(6), Inches(0.5),
                 "TECHNICAL UNIVERSITY OF CHEMNITZ", size=10, bold=True,
                 color=TEXT_LIGHT)
        # CODE GUARDIAN right
        add_text(slide, SLIDE_W - Inches(3.0) - Inches(0.4), Inches(0.22),
                 Inches(3.0), Inches(0.5), "CODE GUARDIAN",
                 size=11, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    if show_title:
        # title bar below header
        title_y = HEADER_H + Inches(0.15)
        add_text(slide, MARGIN, title_y, SLIDE_W - 2 * MARGIN, Inches(0.65),
                 title, size=30, bold=True, color=NAVY)
        if accent_strip:
            strip = add_rect(slide, MARGIN, title_y + Inches(0.65),
                             Inches(0.7), Inches(0.04), fill=ACCENT)
    if footer:
        add_text(slide, MARGIN, SLIDE_H - FOOTER_H,
                 SLIDE_W - 2 * MARGIN - Inches(0.6), FOOTER_H,
                 AUTHOR_FOOTER, size=9, color=TEXT_MUTED,
                 anchor=MSO_ANCHOR.MIDDLE)
        # page number placeholder (filled in post-pass)
        if page_idx is not None:
            add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - FOOTER_H,
                     Inches(0.8), FOOTER_H,
                     f"{page_idx} / __",
                     size=9, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.RIGHT)


def content_box():
    """Top-left + width/height of the main content area below the title."""
    y = HEADER_H + Inches(1.05)
    return (MARGIN, y, SLIDE_W - 2 * MARGIN, SLIDE_H - y - FOOTER_H - Inches(0.1))


def card(slide, x, y, w, h, *, icon=None, title=None, body=None,
         badge=None, accent_color=NAVY, body_size=12, title_size=16,
         body_color=TEXT_DARK):
    add_rect(slide, x, y, w, h, fill=CARD_BG, line=CARD_BORDER, line_pt=0.75,
             rounded=True)
    # accent top strip
    strip = add_rect(slide, x, y, w, Inches(0.08), fill=accent_color)
    strip.line.fill.background()
    pad = Inches(0.22)
    cy = y + Inches(0.20)
    # header line: icon + title + badge
    title_x = x + pad
    title_w = w - 2 * pad
    if icon:
        add_text(slide, title_x, cy, Inches(0.6), Inches(0.45),
                 icon, size=22, color=accent_color, bold=True)
        title_x = title_x + Inches(0.55)
        title_w = title_w - Inches(0.55)
    if badge:
        add_text(slide, x + w - pad - Inches(1.0), cy + Inches(0.05),
                 Inches(1.0), Inches(0.35), badge, size=10, bold=True,
                 color=accent_color, align=PP_ALIGN.RIGHT)
        title_w = title_w - Inches(1.0)
    if title:
        add_text(slide, title_x, cy, title_w, Inches(0.45),
                 title, size=title_size, bold=True, color=NAVY_DARK)
    if body:
        by = cy + Inches(0.55)
        add_text(slide, x + pad, by, w - 2 * pad, h - (by - y) - Inches(0.2),
                 body, size=body_size, color=body_color)


def hero_number_card(slide, x, y, w, h, *, label, value, subtext, color=NAVY):
    add_rect(slide, x, y, w, h, fill=NAVY, line=None, rounded=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), Inches(0.4),
             label, size=12, bold=True, color=RGBColor(0xFC, 0xA3, 0x11))
    add_text(slide, x + Inches(0.25), y + Inches(0.55), w - Inches(0.5), Inches(0.9),
             value, size=36, bold=True, color=TEXT_LIGHT)
    add_text(slide, x + Inches(0.25), y + Inches(1.55), w - Inches(0.5), h - Inches(1.75),
             subtext, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))


def placeholder(slide, x, y, w, h, label, *, sub=None):
    s = add_rect(slide, x, y, w, h, fill=RGBColor(0xEE, 0xF0, 0xF6),
                 line=CARD_BORDER, line_pt=1.0, rounded=True)
    # dashed border effect — approximate with darker line
    add_text(slide, x, y + h/2 - Inches(0.35), w, Inches(0.4),
             "📐  " + label, size=14, bold=True,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        add_text(slide, x, y + h/2, w, Inches(0.4),
                 sub, size=10, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def badge(slide, x, y, w, h, text, *, color=ACCENT):
    add_rect(slide, x, y, w, h, fill=color, line=None, rounded=True)
    add_text(slide, x, y, w, h, text, size=11, bold=True,
             color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 1 — Title
# ============================================================
def slide_01_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    # left accent strip
    add_rect(s, Inches(0.6), Inches(0.7), Inches(0.04), Inches(0.5),
             fill=ACCENT)
    add_text(s, Inches(0.85), Inches(0.65), Inches(8), Inches(0.5),
             "TECHNICAL UNIVERSITY OF CHEMNITZ", size=12, bold=True,
             color=TEXT_LIGHT)
    add_text(s, SLIDE_W - Inches(3.8), Inches(0.55), Inches(3.2), Inches(0.7),
             "CODE GUARDIAN", size=20, bold=True, color=TEXT_LIGHT,
             align=PP_ALIGN.RIGHT)
    # title
    add_text(s, Inches(0.85), Inches(2.0), SLIDE_W - Inches(1.7), Inches(2.0),
             "Privacy-Preserving Source Code\nVulnerability Detection & Repair",
             size=44, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(4.0), SLIDE_W - Inches(1.7), Inches(0.8),
             "Retrieval-Augmented Local LLMs for Visual Studio Code",
             size=20, color=RGBColor(0xC8, 0xD0, 0xE0))
    # presenter columns
    col_y = Inches(5.5)
    col_w = Inches(3.9)
    cols = [
        ("PRESENTED BY", "Md Hafizur Rahman", "Mat. No. 810641  ·  M.Sc. Web Engineering"),
        ("SUPERVISOR  /  EXAMINER", "Abubaker Gaber  (Internal)",
         "Dr.-Ing. Sebastian Heil  (Examiner)\nProf. Dr.-Ing. Martin Gaedke  (Chair)"),
        ("DEFENSE DATE", DATE_STR,
         "Chair of Distributed and Self-Organizing Systems"),
    ]
    for i, (h, name, sub) in enumerate(cols):
        x = Inches(0.85) + i * (col_w + Inches(0.15))
        add_text(s, x, col_y, col_w, Inches(0.3),
                 h, size=9, bold=True, color=ACCENT)
        add_text(s, x, col_y + Inches(0.32), col_w, Inches(0.4),
                 name, size=13, bold=True, color=TEXT_LIGHT)
        add_text(s, x, col_y + Inches(0.72), col_w, Inches(0.8),
                 sub, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))


# ============================================================
# SLIDE 2 — The Problem (illustration placeholder)
# ============================================================
def slide_02_problem_illustration():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "The Developer's Dilemma", page_idx=2)
    x, y, w, h = content_box()
    placeholder(s, x + Inches(1.0), y, w - Inches(2.0), h - Inches(0.3),
                "ILLUSTRATION GOES HERE",
                sub="Developer at a terminal · 3 speech bubbles: \"Can't paste this — customer data\"  ·  \"Semgrep missed it again\"  ·  \"I need a fix, not just a warning\"")


# ============================================================
# SLIDE 3 — What "Privacy" Costs Today
# ============================================================
def slide_03_two_failures():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Today You Pick: Send Code Away — or Get No Help", page_idx=3)
    x, y, w, h = content_box()
    gap = Inches(0.4)
    cw = (w - gap) / 2
    card(s, x, y, cw, h - Inches(0.2),
         icon="❌", title="Cloud LLM (Copilot, ChatGPT, Cursor)",
         body=("Source code crosses the network on every request.\n\n"
               "Blocked outright in finance, healthcare, defence,\n"
               "and air-gapped environments.\n\n"
               "Strong reasoning · zero privacy guarantees."),
         body_size=13, accent_color=RED)
    card(s, x + cw + gap, y, cw, h - Inches(0.2),
         icon="❌", title="Local SAST (Semgrep, CodeQL, ESLint)",
         body=("Stays local. Cannot reason across files or generate fixes.\n\n"
               "Semgrep canonical recall on our JS/TS corpus:  12.68 %\n\n"
               "Strong privacy · poor coverage · no repair."),
         body_size=13, accent_color=RED)


# ============================================================
# SLIDE 4 — The Three-Way Gap
# ============================================================
def slide_04_three_gap():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "No Tool Combines All Three", page_idx=4)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - 2 * gap) / 3
    rows = [
        ("🔒", "Local Inference",
         "Code never leaves the\ndeveloper's machine.\nVerifiable network boundary."),
        ("📚", "Retrieval-Augmented",
         "Grounded in CWE / OWASP\nknowledge.\nNot just pattern matching."),
        ("🛠️", "Developer-Controlled Repair",
         "Structured fix returned.\nDeveloper applies via Quick Fix.\nNo silent rewrites."),
    ]
    for i, (icon, t, b) in enumerate(rows):
        cx = x + i * (cw + gap)
        card(s, cx, y, cw, Inches(3.2),
             icon=icon, title=t, body=b, body_size=13, title_size=18,
             accent_color=ACCENT)
    # callout below
    co_y = y + Inches(3.5)
    add_rect(s, x, co_y, w, Inches(0.9), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), co_y, w - Inches(0.6), Inches(0.9),
             "This thesis builds and evaluates a tool that delivers all three.",
             size=18, bold=True, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — Research Questions & Requirements
# ============================================================
def slide_05_rqs_reqs():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Research Questions and Requirements", page_idx=5)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    lw = w * 0.40 - gap / 2
    rw = w * 0.60 - gap / 2

    # LEFT — RQs
    add_rect(s, x, y, lw, Inches(3.6), fill=CARD_BG,
             line=CARD_BORDER, line_pt=0.75, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(0.15), lw - Inches(0.6), Inches(0.4),
             "RESEARCH QUESTIONS", size=11, bold=True, color=ACCENT)
    rqs = [
        ("RQ1", "Can a local LLM on consumer hardware reach acceptable detection quality under a strict privacy envelope?"),
        ("RQ2", "Does grounding the LLM in retrieved CWE / OWASP knowledge (RAG) improve detection?"),
        ("RQ3", "Is in-IDE latency acceptable for developer workflows?"),
    ]
    row_y = y + Inches(0.65)
    for tag, q in rqs:
        add_text(s, x + Inches(0.3), row_y, Inches(0.7), Inches(0.4),
                 tag, size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(1.05), row_y, lw - Inches(1.35), Inches(1.0),
                 q, size=11, color=TEXT_DARK)
        row_y += Inches(0.95)

    # RIGHT — R-cards
    add_rect(s, x + lw + gap, y, rw, Inches(3.6), fill=CARD_BG,
             line=CARD_BORDER, line_pt=0.75, rounded=True)
    add_text(s, x + lw + gap + Inches(0.3), y + Inches(0.15), rw - Inches(0.6),
             Inches(0.4),
             "FIVE REQUIREMENTS", size=11, bold=True, color=ACCENT)
    rs = [
        ("R1", "Accuracy", "precision, recall, F1"),
        ("R2", "Consistency", "JSON parse, inter-run"),
        ("R3", "Repair Quality", "manual + auto-applicable"),
        ("R4", "Usability", "latency, IDE integration"),
        ("R5", "Privacy", "zero exfiltration, signed corpus"),
    ]
    rw_inner = (rw - Inches(0.6) - Inches(0.3)) / 2
    rx = x + lw + gap + Inches(0.3)
    ry = y + Inches(0.65)
    for i, (tag, lbl, sub) in enumerate(rs):
        col = i % 2
        row = i // 2
        cx = rx + col * (rw_inner + Inches(0.3))
        cy = ry + row * Inches(0.85)
        add_rect(s, cx, cy, rw_inner, Inches(0.75),
                 fill=LIGHT_BG, line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.15), cy + Inches(0.06), Inches(0.5), Inches(0.35),
                 tag, size=12, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.75), cy + Inches(0.06),
                 rw_inner - Inches(0.85), Inches(0.35),
                 lbl, size=13, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.75), cy + Inches(0.36),
                 rw_inner - Inches(0.85), Inches(0.35),
                 sub, size=9, color=TEXT_MUTED)

    # callout
    add_rect(s, x, y + Inches(3.75), w, Inches(0.55), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(3.75), w - Inches(0.6), Inches(0.55),
             "R1–R5 use fixed four-level threshold scales set in Chapter 2 — before evaluation.",
             size=12, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — Existing Approaches
# ============================================================
def slide_06_existing_approaches():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Three Existing Approaches", page_idx=6)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 2 * gap) / 3
    cols = [
        ("CLOUD LLM ASSISTANTS",
         [("GitHub Copilot", "Microsoft · 2021"),
          ("ChatGPT / GPT-4", "OpenAI · 2022"),
          ("Cursor", "Anysphere · 2023"),
          ("Claude Code", "Anthropic · 2024")],
         RED),
        ("LOCAL SAST",
         [("Semgrep", "Rule-based · OSS"),
          ("CodeQL", "GitHub · taint analysis"),
          ("ESLint security plugins", "Lint-level"),
          ("Bandit / SpotBugs", "Other languages")],
         ACCENT),
        ("ACADEMIC LLM-FOR-SEC",
         [("IRIS + GPT-4", "Li et al. · ICLR 2025"),
          ("SecRepair", "Islam et al. · NDSS 2024"),
          ("LLMSecGuard", "Kavian et al. · EASE 2024"),
          ("RESCUE", "Shi & Zhang · 2025")],
         NAVY),
    ]
    for i, (hdr, items, color) in enumerate(cols):
        cx = x + i * (cw + gap)
        # header band
        add_rect(s, cx, y, cw, Inches(0.55), fill=NAVY, rounded=True)
        add_text(s, cx, y, cw, Inches(0.55), hdr, size=11, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rowy = y + Inches(0.75)
        for name, sub in items:
            add_rect(s, cx, rowy, cw, Inches(0.7), fill=CARD_BG,
                     line=CARD_BORDER, rounded=True)
            add_text(s, cx + Inches(0.2), rowy + Inches(0.08), cw - Inches(0.4),
                     Inches(0.3), name, size=12, bold=True, color=NAVY_DARK)
            add_text(s, cx + Inches(0.2), rowy + Inches(0.38), cw - Inches(0.4),
                     Inches(0.3), sub, size=10, color=TEXT_MUTED, font="Calibri")
            rowy += Inches(0.8)


# ============================================================
# SLIDE 7 — Critical Gap Matrix
# ============================================================
def slide_07_gap_matrix():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "No Tool Satisfies R1 + R3 + R5 Together", page_idx=7)
    x, y, w, h = content_box()

    rows = [
        ("GitHub Copilot",   ["✅","🟡","✅","✅","❌"]),
        ("Semgrep",          ["🟡","✅","❌","✅","✅"]),
        ("CodeQL",           ["🟡","✅","❌","🟡","✅"]),
        ("IRIS (GPT-4)",     ["🟡","🟡","❌","❌","❌"]),
        ("SecRepair",        ["🟡","🟡","🟡","❌","🟡"]),
        ("Code Guardian",    ["✅","✅","✅","✅","✅"]),
    ]
    header = ["Solution / Approach", "R1\nACCURACY", "R2\nCONSISTENCY",
              "R3\nREPAIR", "R4\nUSABILITY", "R5\nPRIVACY"]

    table_w = w * 0.75
    legend_w = w - table_w - Inches(0.3)
    col_w = [Inches(2.3)] + [(table_w - Inches(2.3)) / 5] * 5
    row_h = Inches(0.5)

    # header row
    cx = x
    for i, hd in enumerate(header):
        add_rect(s, cx, y, col_w[i], Inches(0.65), fill=NAVY)
        add_text(s, cx, y, col_w[i], Inches(0.65), hd,
                 size=10, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    ry = y + Inches(0.65)
    for r_i, (name, cells) in enumerate(rows):
        is_self = name == "Code Guardian"
        bg = RGBColor(0xE9, 0xF3, 0xEA) if is_self else (LIGHT_BG if r_i % 2 == 0 else CARD_BG)
        cx = x
        add_rect(s, cx, ry, col_w[0], row_h, fill=bg, line=CARD_BORDER)
        add_text(s, cx + Inches(0.15), ry, col_w[0] - Inches(0.2), row_h,
                 name, size=12, bold=is_self, color=NAVY_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[0]
        for c_i, cell in enumerate(cells):
            add_rect(s, cx, ry, col_w[c_i+1], row_h, fill=bg, line=CARD_BORDER)
            add_text(s, cx, ry, col_w[c_i+1], row_h, cell,
                     size=16, color=NAVY_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c_i+1]
        ry += row_h

    # legend + callout
    lx = x + table_w + Inches(0.3)
    ly = y
    add_rect(s, lx, ly, legend_w, Inches(1.4),
             fill=CARD_BG, line=CARD_BORDER, rounded=True)
    add_text(s, lx + Inches(0.15), ly + Inches(0.1), legend_w - Inches(0.3),
             Inches(0.3), "LEGEND", size=10, bold=True, color=ACCENT)
    add_text(s, lx + Inches(0.15), ly + Inches(0.4), legend_w - Inches(0.3),
             Inches(1.0),
             "✅  High / Satisfied\n🟡  Partial\n❌  Not satisfied",
             size=11, color=TEXT_DARK)

    co_y = ly + Inches(1.6)
    add_rect(s, lx, co_y, legend_w, Inches(2.2), fill=NAVY, rounded=True)
    add_text(s, lx + Inches(0.2), co_y + Inches(0.15), legend_w - Inches(0.4),
             Inches(0.4),
             "THE GAP", size=11, bold=True, color=ACCENT)
    add_text(s, lx + Inches(0.2), co_y + Inches(0.5), legend_w - Inches(0.4),
             Inches(1.6),
             "No existing tool delivers privacy + repair + IDE-grade latency together.",
             size=14, bold=True, color=TEXT_LIGHT)


# ============================================================
# SLIDE 8 — Four Developer Workflows
# ============================================================
def slide_08_workflows():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "One Pipeline, Four Developer Workflows", page_idx=8)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 3 * gap) / 4
    rows = [
        ("⚡", "Real-time Diagnostics", "Inline squiggles per function as you save."),
        ("📂", "On-demand File Scan", "Right-click → scan the current file."),
        ("💬", "Interactive Q&A", "Ask about a finding in a side panel."),
        ("🗂️", "Workspace Audit", "Batch scan the whole project."),
    ]
    for i, (icon, t, b) in enumerate(rows):
        cx = x + i * (cw + gap)
        card(s, cx, y, cw, Inches(2.8),
             icon=icon, title=t, body=b, body_size=12, title_size=14,
             accent_color=ACCENT)

    # below: pipeline strip
    py = y + Inches(3.1)
    pipe_h = Inches(1.6)
    add_rect(s, x, py, w, pipe_h, fill=CARD_BG,
             line=CARD_BORDER, rounded=True)
    add_text(s, x + Inches(0.3), py + Inches(0.12), w - Inches(0.6), Inches(0.3),
             "ALL FOUR WORKFLOWS HIT THE SAME PIPELINE", size=10, bold=True,
             color=ACCENT)
    # boxes
    nodes = ["INPUT", "STAGE 1\nDetection (LLM + RAG)", "Consensus\nfilter",
             "STAGE 2\nRepair (LLM)", "QUICK\nFIX"]
    bx = x + Inches(0.4)
    by = py + Inches(0.55)
    bw = (w - Inches(0.8) - 4 * Inches(0.25)) / 5
    bh = Inches(0.85)
    for i, n in enumerate(nodes):
        cx = bx + i * (bw + Inches(0.25))
        fill = NAVY if i in (1, 3) else RGBColor(0xE6, 0xEB, 0xF5)
        col = TEXT_LIGHT if i in (1, 3) else NAVY
        add_rect(s, cx, by, bw, bh, fill=fill, rounded=True)
        add_text(s, cx, by, bw, bh, n, size=11, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(nodes) - 1:
            ax = cx + bw + Inches(0.02)
            add_text(s, ax, by, Inches(0.21), bh, "→", size=18,
                     color=ACCENT, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 9 — Two-stage pipeline detail
# ============================================================
def slide_09_two_stage():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Two-Stage Local Pipeline", page_idx=9)
    x, y, w, h = content_box()
    gap = Inches(0.4)
    cw = (w - gap) / 2
    card(s, x, y, cw, Inches(3.6),
         icon="🔍", title="Stage 1 — Detection",
         body=("• JSON-mode LLM call  ·  schema-pinned output\n"
               "• Optional RAG context  ·  top-k CWE / OWASP snippets\n"
               "• Three deterministic runs  ·  seed = 42, temperature = 0\n"
               "• Inter-run confidence gate  ·  ≥ 2-of-3 agreement (0.67)\n"
               "• Returns:  category, severity, line range, evidence"),
         body_size=12, accent_color=NAVY)
    card(s, x + cw + gap, y, cw, Inches(3.6),
         icon="🛠️", title="Stage 2 — Repair",
         body=("• Separate structured call  ·  invoked only when detected\n"
               "• Returns  { code, language }  ·  18 explicit abstentions allowed\n"
               "• Validator parses output with @babel/parser  ·  rejects prose\n"
               "• Surfaced as a VS Code Quick Fix\n"
               "• Developer applies the patch  ·  no silent rewrites"),
         body_size=12, accent_color=ACCENT)
    # bottom callout
    by = y + Inches(3.8)
    add_rect(s, x, by, w, Inches(0.6), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by, w - Inches(0.6), Inches(0.6),
             "Privacy by construction  ·  Ollama bound to loopback  ·  embeddings local  ·  no telemetry",
             size=13, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — Live: VS Code extension screenshot
# ============================================================
def slide_10_screenshot():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Live in VS Code", page_idx=10)
    x, y, w, h = content_box()
    placeholder(s, x, y, w - Inches(2.2), h - Inches(0.2),
                "SCREENSHOT GOES HERE",
                sub="VS Code with: a JS/TS file open · red squiggle on a vulnerable line · Code Guardian side panel showing the finding · Quick Fix menu visible")
    # QR placeholder column
    qx = x + w - Inches(2.0)
    placeholder(s, qx, y, Inches(2.0), Inches(2.0), "QR CODE",
                sub="Demo video / GitHub release")
    # caption
    add_text(s, qx, y + Inches(2.1), Inches(2.0), Inches(0.4),
             "SCAN  ·  Demo", size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — System Architecture
# ============================================================
def slide_11_architecture():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "System Architecture", page_idx=11)
    x, y, w, h = content_box()
    placeholder(s, x, y, w, h - Inches(0.2),
                "ARCHITECTURE DIAGRAM",
                sub="VS Code Extension Host → Ollama (loopback) ↔ RAG store (HNSWlib) ↔ Vulnerability Data Manager → NVD / OWASP / CWE feeds  ·  Dashed rectangle = the privacy boundary (everything that stays on the laptop)")


# ============================================================
# SLIDE 12 — Privacy by Construction
# ============================================================
def slide_12_privacy_cards():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Privacy by Construction", page_idx=12)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(4.0) - gap) / 2
    items = [
        ("🔒", "Loopback-only Ollama",
         "Bound to 127.0.0.1. Network isolation arm of the threat model — empirically verified."),
        ("🧬", "Signed RAG Corpus",
         "Ed25519 manifest verified at load. Provenance arm of the threat model — corpus cannot be silently rewritten."),
        ("📦", "Pinned Container",
         "Node 20.19.0-alpine · npm ci · seed = 42 · byte-identical runs."),
        ("🛡️", "Prompt-Injection Harness",
         "12 cases of crafted prompts attempting exfiltration. leakFreeRate = 100 %."),
    ]
    for i, (icon, t, b) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=icon, title=t, body=b,
             body_size=12, title_size=15, accent_color=ACCENT)


# ============================================================
# SLIDE 13 — Evaluation Setup
# ============================================================
def slide_13_eval_setup():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Evaluation Setup", page_idx=13)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - 2 * gap) / 3
    items = [
        ("📚", "Corpus",
         "Curated JS / TS  ·  101 cases\n(71 vulnerable + 30 secure)\n14 CWE categories\n\nExternal:  15 cases from NodeGoat,\nJuice Shop, 3 named CVEs\n\nWhole-project scan:  OWASP NodeGoat"),
        ("🧪", "Configurations",
         "5 Ollama models × {LLM-only, LLM+RAG}\n= 10 configurations\n\n+ 3 SAST baselines\n(Semgrep, CodeQL, ESLint)\n\n3 runs per sample  ·  deterministic\n→  303 invocations per config"),
        ("📐", "Statistical Discipline",
         "Held-out 71 / 30 test split\n(threshold tuning off the test set)\n\nMcNemar with Bonferroni\nacross paired model comparisons\n\nExact-binomial CIs for small-n rates"),
    ]
    for i, (icon, t, b) in enumerate(items):
        cx = x + i * (cw + gap)
        card(s, cx, y, cw, Inches(4.3),
             icon=icon, title=t, body=b, body_size=11, title_size=15,
             accent_color=NAVY)


# ============================================================
# SLIDE 14 — Hero numbers
# ============================================================
def slide_14_hero():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Headline Numbers", page_idx=14)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2
    ch = (Inches(4.0) - gap) / 2
    cards = [
        ("R1 — DETECTION",  "71.43 %",
         "F1  ·  qwen3:8b + RAG  ·  precision 72.46 %  ·  recall 70.42 %\nWith confidence gate ≥ 0.67  →  74.07 % F1, 78.13 % precision"),
        ("R3 — REPAIR",     "90.32 %",
         "Auto-applicable rate (n = 279)\nManual review: 89.5 % semantic, 68 % combined (n = 25)"),
        ("R4 — LATENCY",    "2.2 s",
         "Median stage-1  ·  p95 = 4.3 s  ·  comfortably inside the on-demand band (≤ 5 s)"),
        ("R5 — PRIVACY",    "100 %",
         "leakFreeRate on 12-case prompt-injection harness\nZero non-loopback transmission across the corpus run"),
    ]
    for i, (lbl, val, sub) in enumerate(cards):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        hero_number_card(s, cx, cy, cw, ch, label=lbl, value=val, subtext=sub)
    # bottom strip
    by = y + 2 * ch + gap + Inches(0.05)
    add_text(s, x, by, w, Inches(0.35),
             "Each number detailed on the next four slides.",
             size=12, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 — R1 detection accuracy chart
# ============================================================
def slide_15_r1_accuracy():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R1 — Detection Accuracy", page_idx=15)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R1", color=ACCENT)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    # top: numbers card
    add_rect(s, x, y, w, Inches(1.4), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), y + Inches(0.15), w - Inches(0.6), Inches(0.4),
             "qwen3:8b + RAG  ·  canonical taxonomy  ·  full 101-case run",
             size=11, bold=True, color=ACCENT)
    # three big stats
    stats = [
        ("71.43 %", "F1"),
        ("72.46 %", "Precision"),
        ("70.42 %", "Recall"),
        ("10.0 %", "FPR"),
    ]
    sw = (w - Inches(0.4)) / 4
    for i, (v, lbl) in enumerate(stats):
        sx = x + Inches(0.2) + i * sw
        add_text(s, sx, y + Inches(0.55), sw, Inches(0.5),
                 v, size=24, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(s, sx, y + Inches(1.05), sw, Inches(0.3),
                 lbl, size=11, color=RGBColor(0xC8, 0xD0, 0xE0),
                 align=PP_ALIGN.CENTER)

    # below: F1 vs SAST baselines bar chart placeholder
    chart_y = y + Inches(1.6)
    placeholder(s, x, chart_y, w * 0.62, h - Inches(1.7),
                "BAR CHART  ·  F1: Code Guardian vs SAST baselines",
                sub="Code Guardian 71.43 %  ·  Semgrep ~14 %  ·  CodeQL ~X  ·  ESLint-security ~X")

    # right: confidence gate + held-out side card
    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    card(s, rx, chart_y, rw, h - Inches(1.7),
         icon="📈", title="Refinements",
         body=("With confidence gate ≥ 0.67:\n"
               "F1 = 74.07 %  ·  precision = 78.13 %\n\n"
               "Held-out 71 / 30 test split:\n"
               "F1 = 61.11 %  ·  FPR = 0 %\n\n"
               "Beats every SAST baseline\non recall-driven F1."),
         body_size=11, title_size=14, accent_color=ACCENT)


# ============================================================
# SLIDE 16 — RAG ablation
# ============================================================
def slide_16_rag_ablation():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "RAG Is Not a Uniform Win", page_idx=16)
    badge(s, SLIDE_W - Inches(2.4), Inches(0.25), Inches(1.8), Inches(0.4),
          "R1  ·  RQ2", color=ACCENT)
    x, y, w, h = content_box()

    # chart placeholder
    placeholder(s, x, y, w * 0.62, h - Inches(0.2),
                "BAR CHART  ·  ΔF1 from no-RAG → RAG per model",
                sub="qwen3:8b +3.91  ·  gemma3:1b ≈ 0  ·  gemma3:4b ≈ 0  ·  qwen3:4b −5.22  ·  codellama −29.89")

    # right: take-away
    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    add_rect(s, rx, y, rw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.2), y + Inches(0.15), rw - Inches(0.4), Inches(0.4),
             "TAKEAWAY", size=11, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.2), y + Inches(0.55), rw - Inches(0.4), Inches(2.5),
             "RAG can help, harm,\nor do nothing —\ndepending on the model.",
             size=20, bold=True, color=TEXT_LIGHT)
    add_text(s, rx + Inches(0.2), y + Inches(2.9), rw - Inches(0.4), Inches(1.8),
             "Only codellama's −29.89 F1 survives Bonferroni at α = 0.01\n(McNemar p = 0.0013).\n\nDesign implication:  treat RAG as a per-model configuration choice, not a default.",
             size=11, color=RGBColor(0xD8, 0xDE, 0xEE))


# ============================================================
# SLIDE 17 — R3 repair quality
# ============================================================
def slide_17_r3_repair():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R3 — Repair Quality", page_idx=17)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R3", color=ACCENT)
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2

    # left card — manual
    add_rect(s, x, y, cw, Inches(3.4), fill=CARD_BG,
             line=CARD_BORDER, rounded=True)
    add_rect(s, x, y, cw, Inches(0.08), fill=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
             "👤  Manual review  ·  n = 25", size=14, bold=True, color=NAVY_DARK)
    add_text(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.5), Inches(0.4),
             "GENEROUS BAR", size=10, bold=True, color=TEXT_MUTED)
    # numbers
    add_text(s, x + Inches(0.25), y + Inches(1.15), cw - Inches(0.5), Inches(0.7),
             "89.5 %", size=42, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(1.95), cw - Inches(0.5), Inches(0.3),
             "semantic correctness  ·  fix addresses the CWE",
             size=11, color=TEXT_DARK)
    add_text(s, x + Inches(0.25), y + Inches(2.35), cw - Inches(0.5), Inches(0.4),
             "68 %", size=24, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(2.8), cw - Inches(0.5), Inches(0.4),
             "combined correctness  ·  semantic + strategy + executable",
             size=11, color=TEXT_MUTED)

    # right — auto
    rx = x + cw + gap
    add_rect(s, rx, y, cw, Inches(3.4), fill=CARD_BG,
             line=CARD_BORDER, rounded=True)
    add_rect(s, rx, y, cw, Inches(0.08), fill=ACCENT)
    add_text(s, rx + Inches(0.25), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
             "⚙️  Auto-applicable rate  ·  fully automated", size=14, bold=True,
             color=NAVY_DARK)
    add_text(s, rx + Inches(0.25), y + Inches(0.75), cw - Inches(0.5), Inches(0.4),
             "STRICT BAR  ·  validator parses with @babel/parser",
             size=10, bold=True, color=TEXT_MUTED)
    add_text(s, rx + Inches(0.25), y + Inches(1.15), cw - Inches(0.5), Inches(0.7),
             "90.32 %", size=42, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.25), y + Inches(1.95), cw - Inches(0.5), Inches(0.3),
             "on issued fixes  (252 / 279)",
             size=11, color=TEXT_DARK)
    add_text(s, rx + Inches(0.25), y + Inches(2.35), cw - Inches(0.5), Inches(0.4),
             "84.85 %", size=24, bold=True, color=ACCENT)
    add_text(s, rx + Inches(0.25), y + Inches(2.8), cw - Inches(0.5), Inches(0.4),
             "across all stage-2 calls  (no-fix abstentions count as failures)",
             size=11, color=TEXT_MUTED)

    # bottom callout
    by = y + Inches(3.55)
    add_rect(s, x, by, w, Inches(0.85), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6), Inches(0.3),
             "TWO-TIER METRIC", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.35), w - Inches(0.6), Inches(0.5),
             "Manual = decision-support reading  ·  Auto-applicable = deployment reading.  Stage 2 returns { code, language } + 18 explicit abstentions.",
             size=12, color=TEXT_LIGHT)


# ============================================================
# SLIDE 18 — R4 latency bands
# ============================================================
def slide_18_r4_latency():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R4 — Usability: Latency Bands", page_idx=18)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R4", color=ACCENT)
    x, y, w, h = content_box()
    # left — chart
    placeholder(s, x, y, w * 0.62, h - Inches(0.2),
                "HORIZONTAL BAR CHART  ·  median latency per model",
                sub="Vertical dashed lines: 500 ms (real-time) · 1.5 s (interactive) · 5 s (on-demand)")
    # right — band cards
    rx = x + w * 0.62 + Inches(0.2)
    rw = w - w * 0.62 - Inches(0.2)
    bands = [
        ("⚡  Real-time  ≤ 500 ms", "SAST baselines only", RED),
        ("🖱️  Interactive  ≤ 1.5 s", "gemma3:1b + RAG  (1,019 ms)\nqwen3:4b  (1,328 ms)", ACCENT),
        ("🧠  On-demand  ≤ 5 s", "qwen3:8b + RAG  (2,216 ms)  ←  headline\np95 = 4,327 ms", GREEN),
    ]
    by = y
    bh = (h - Inches(0.4) - 2 * Inches(0.2)) / 3
    for hdr, body, col in bands:
        add_rect(s, rx, by, rw, bh, fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_rect(s, rx, by, Inches(0.1), bh, fill=col)
        add_text(s, rx + Inches(0.2), by + Inches(0.1), rw - Inches(0.3),
                 Inches(0.4), hdr, size=12, bold=True, color=NAVY_DARK)
        add_text(s, rx + Inches(0.2), by + Inches(0.45), rw - Inches(0.3),
                 bh - Inches(0.5), body, size=11, color=TEXT_DARK)
        by += bh + Inches(0.2)


# ============================================================
# SLIDE 19 — R5 privacy
# ============================================================
def slide_19_r5_privacy():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "R5 — Privacy: Empirically Verified", page_idx=19)
    badge(s, SLIDE_W - Inches(1.3), Inches(0.25), Inches(0.7), Inches(0.4),
          "R5", color=ACCENT)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(3.8) - gap) / 2
    items = [
        ("✅", "Zero non-loopback traffic",
         "Across the full 303-invocation corpus run."),
        ("✅", "leakFreeRate = 100 %",
         "12-case prompt-injection harness  ·  no exfiltration on any case."),
        ("✅", "Ed25519-signed manifest",
         "RAG corpus provenance verified at load."),
        ("✅", "Reproducible",
         "Node 20.19.0-alpine  ·  npm ci  ·  seed = 42."),
    ]
    for i, (icon, t, b) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=icon, title=t, body=b,
             body_size=12, title_size=15, accent_color=GREEN)
    # footer info
    by = y + 2 * ch + gap + Inches(0.1)
    add_rect(s, x, by, w, Inches(0.55), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by, w - Inches(0.6), Inches(0.55),
             "Resource footprint:  harness 3 ms median CPU, 87 MB RSS  ·  Ollama 6–8 GB RAM for 8B-q4.",
             size=11, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 20 — What We Learned
# ============================================================
def slide_20_takeaways():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "What We Learned", page_idx=20)
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(4.3) - gap) / 2
    items = [
        ("🏆", "A local 8B model is enough",
         "qwen3:8b + RAG reaches ~71 % F1 on JS/TS — well inside the on-demand latency band on a developer laptop. The privacy envelope holds without paying for it in quality."),
        ("📚", "RAG is model-dependent",
         "The largest paired effect in the dataset is a RAG-induced degradation (codellama −29.89 F1). Treat retrieval as a per-model design choice, not a default."),
        ("🛠️", "Repair is the real product gap",
         "90 % auto-applicable means the developer can act, not just be warned. This is where SAST falls flat — they detect, they don't fix."),
        ("🔒", "Privacy as a verifiable boundary",
         "Loopback + signed corpus + container pin makes privacy testable, not aspirational. leakFreeRate = 100 % isn't a slogan — it's a measurement."),
    ]
    for i, (icon, t, b) in enumerate(items):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        card(s, cx, cy, cw, ch, icon=icon, title=t, body=b,
             body_size=11, title_size=15, accent_color=NAVY)


# ============================================================
# SLIDE 21 — Limitations
# ============================================================
def slide_21_limitations():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Limitations  ·  Honest", page_idx=21)
    x, y, w, h = content_box()
    items = [
        ("👤", "Single-reviewer manual review (n = 25)",
         "Inter-rater agreement on a 10-sample subset = future work.  Headline R3 metric is the fully objective auto-applicable rate."),
        ("🎯", "Curated 101-case corpus",
         "Selection bias mitigated by held-out 71 / 30 split, 15-case external set (NodeGoat, Juice Shop, 3 CVEs), and whole-project NodeGoat run."),
        ("🌐", "One real-world project end-to-end",
         "NodeGoat  ·  multi-project validation is future work."),
        ("🔁", "Consensus filter inert under deterministic decoding",
         "Three runs are byte-identical with seed = 42.  Needs seed rotation to produce a real consistency signal."),
        ("☁️", "No cloud-LLM baseline",
         "Excluded by the privacy threat model  ·  floor provided by 3 SAST baselines."),
    ]
    rh = (h - Inches(0.2)) / len(items)
    for i, (icon, t, b) in enumerate(items):
        ry = y + i * rh
        add_rect(s, x, ry, w, rh - Inches(0.1), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_rect(s, x, ry, Inches(0.1), rh - Inches(0.1), fill=RED)
        add_text(s, x + Inches(0.2), ry + Inches(0.1), Inches(0.6), rh - Inches(0.2),
                 icon, size=20, color=RED, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), ry + Inches(0.08), Inches(4.5), Inches(0.4),
                 t, size=13, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.85), ry + Inches(0.45), w - Inches(1.0),
                 rh - Inches(0.5), b, size=11, color=TEXT_DARK)


# ============================================================
# SLIDE 22 — Future Work
# ============================================================
def slide_22_future_work():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Future Work", page_idx=22)
    x, y, w, h = content_box()
    items = [
        ("🎲", "Seed rotation",
         "Genuine consensus signal beyond the deterministic-decoding floor."),
        ("👥", "Inter-rater agreement",
         "Second reviewer on the 25-sample subset for κ statistic."),
        ("🏗️", "Multi-project real-world validation",
         "Beyond NodeGoat  ·  measured recall on diverse production codebases."),
        ("⚡", "Streaming partial results",
         "Break the real-time band by surfacing first-token findings."),
        ("🎯", "Adaptive RAG gating",
         "Per-model RAG enable/disable driven by slide 16's finding."),
    ]
    rh = (h - Inches(0.2)) / len(items)
    for i, (icon, t, b) in enumerate(items):
        ry = y + i * rh
        add_rect(s, x, ry, w, rh - Inches(0.1), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_rect(s, x, ry, Inches(0.1), rh - Inches(0.1), fill=ACCENT)
        add_text(s, x + Inches(0.2), ry + Inches(0.1), Inches(0.6), rh - Inches(0.2),
                 icon, size=20, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), ry + Inches(0.08), Inches(4.5), Inches(0.4),
                 t, size=13, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.85), ry + Inches(0.45), w - Inches(1.0),
                 rh - Inches(0.5), b, size=11, color=TEXT_DARK)


# ============================================================
# SLIDE 23 — Contributions
# ============================================================
def slide_23_contributions():
    s = prs.slides.add_slide(BLANK)
    add_slide_chrome(s, "Contributions of This Thesis", page_idx=23)
    x, y, w, h = content_box()
    items = [
        ("🛠️", "Working VS Code extension",
         "Local LLM + RAG + signed corpus + pinned container — published to the Marketplace."),
        ("🔗", "Two-stage pipeline with structured contracts",
         "Detection (JSON-mode) + Repair ({ code, language }) — schema-pinned at every boundary."),
        ("📊", "Empirical study",
         "5 models × 2 modes + 3 SAST baselines · 303 invocations per config · CWE-mapped JS/TS corpus."),
        ("🔁", "Reproducibility infrastructure",
         "Byte-identical runs, Ed25519-signed manifest, pinned Node 20.19.0-alpine container."),
        ("📐", "Statistical discipline",
         "Bonferroni-corrected McNemar, exact-binomial CIs, held-out 71 / 30 split."),
    ]
    rh = (h - Inches(0.2)) / len(items)
    for i, (icon, t, b) in enumerate(items):
        ry = y + i * rh
        add_rect(s, x, ry, w, rh - Inches(0.1), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_rect(s, x, ry, Inches(0.1), rh - Inches(0.1), fill=GREEN)
        add_text(s, x + Inches(0.2), ry + Inches(0.1), Inches(0.6), rh - Inches(0.2),
                 icon, size=20, color=GREEN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), ry + Inches(0.08), Inches(4.5), Inches(0.4),
                 t, size=13, bold=True, color=NAVY_DARK)
        add_text(s, x + Inches(0.85), ry + Inches(0.45), w - Inches(1.0),
                 rh - Inches(0.5), b, size=11, color=TEXT_DARK)


# ============================================================
# SLIDE 24 — Thank You / Q&A
# ============================================================
def slide_24_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, Inches(0.6), Inches(0.7), Inches(0.04), Inches(0.5),
             fill=ACCENT)
    add_text(s, Inches(0.85), Inches(0.65), Inches(8), Inches(0.5),
             "TECHNICAL UNIVERSITY OF CHEMNITZ", size=12, bold=True,
             color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(2.2), SLIDE_W - Inches(1.7), Inches(1.4),
             "Thank You", size=72, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), Inches(3.7), SLIDE_W - Inches(1.7), Inches(0.8),
             "Questions & Discussion", size=30, color=RGBColor(0xC8, 0xD0, 0xE0))

    # headline strip
    strip_y = Inches(5.4)
    add_rect(s, Inches(0.85), strip_y, SLIDE_W - Inches(1.7), Inches(0.85),
             fill=NAVY_DARK, rounded=True)
    add_text(s, Inches(1.0), strip_y + Inches(0.1), SLIDE_W - Inches(2.0),
             Inches(0.3), "HEADLINE", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(1.0), strip_y + Inches(0.35), SLIDE_W - Inches(2.0),
             Inches(0.4),
             "F1 71.43 %   ·   repair 90.32 % auto-applicable   ·   latency 2.2 s median   ·   leakFreeRate 100 %",
             size=13, bold=True, color=TEXT_LIGHT)

    # bottom info
    by = Inches(6.55)
    add_text(s, Inches(0.85), by, Inches(4.5), Inches(0.3),
             "Md Hafizur Rahman", size=12, bold=True, color=TEXT_LIGHT)
    add_text(s, Inches(0.85), by + Inches(0.28), Inches(4.5), Inches(0.3),
             "M.Sc. Web Engineering  ·  TU Chemnitz",
             size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
    add_text(s, Inches(9.5), by, Inches(3.5), Inches(0.3),
             "DEFENSE  ·  " + DATE_STR, size=10, bold=True, color=ACCENT,
             align=PP_ALIGN.RIGHT)


# ============================================================
# APPENDIX
# ============================================================
def appendix_header(slide, title):
    add_slide_chrome(slide, title, page_idx=None)
    badge(slide, SLIDE_W - Inches(1.8), Inches(0.25), Inches(1.2), Inches(0.4),
          "APPENDIX", color=ACCENT)


def slide_A1_full_results():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A1  ·  Full Per-Model Results")
    x, y, w, h = content_box()
    headers = ["Configuration", "F1", "Precision", "Recall", "FPR",
               "Latency p50", "Latency p95", "VRAM"]
    rows = [
        ("qwen3:8b + RAG", "71.43", "72.46", "70.42", "10.0", "2,216 ms", "4,327 ms", "~6 GB"),
        ("qwen3:8b (no RAG)", "67.52", "—", "—", "—", "—", "—", "~6 GB"),
        ("qwen3:4b + RAG", "—", "—", "—", "—", "1,328 ms", "—", "~4 GB"),
        ("qwen3:4b (no RAG)", "—", "—", "—", "—", "—", "—", "~4 GB"),
        ("gemma3:4b + RAG", "—", "—", "—", "—", "—", "—", "~4 GB"),
        ("gemma3:4b (no RAG)", "—", "—", "—", "—", "—", "—", "~4 GB"),
        ("gemma3:1b + RAG", "—", "—", "—", "—", "1,019 ms", "—", "~2 GB"),
        ("gemma3:1b (no RAG)", "—", "—", "—", "—", "—", "—", "~2 GB"),
        ("codellama + RAG", "23.17", "—", "53.52", "—", "3,529 ms", "—", "~6 GB"),
        ("codellama (no RAG)", "53.06", "—", "73.24", "—", "2,024 ms", "—", "~6 GB"),
        ("Semgrep", "~14", "—", "12.68", "low", "<500 ms", "—", "n/a"),
        ("CodeQL", "—", "—", "—", "—", "<500 ms", "—", "n/a"),
        ("ESLint-security", "—", "—", "—", "—", "<500 ms", "—", "n/a"),
    ]
    col_w = [Inches(2.2), Inches(0.9), Inches(1.1), Inches(1.0), Inches(0.9),
             Inches(1.4), Inches(1.4), Inches(1.0)]
    row_h = Inches(0.32)
    cx = x
    for i, hd in enumerate(headers):
        add_rect(s, cx, y, col_w[i], Inches(0.4), fill=NAVY)
        add_text(s, cx, y, col_w[i], Inches(0.4), hd, size=10, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]
    ry = y + Inches(0.4)
    for r_i, row in enumerate(rows):
        is_head = row[0] == "qwen3:8b + RAG"
        bg = RGBColor(0xE9, 0xF3, 0xEA) if is_head else (LIGHT_BG if r_i % 2 == 0 else CARD_BG)
        cx = x
        for c_i, val in enumerate(row):
            add_rect(s, cx, ry, col_w[c_i], row_h, fill=bg, line=CARD_BORDER)
            add_text(s, cx, ry, col_w[c_i], row_h, val, size=9.5,
                     bold=(c_i == 0 or is_head),
                     color=NAVY_DARK,
                     align=PP_ALIGN.CENTER if c_i > 0 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c_i]
        ry += row_h
    # caption
    add_text(s, x, ry + Inches(0.05), w, Inches(0.3),
             "Fill in `—` cells from evaluation/results/evaluation-phase1-b3-2026-04-28-qwen3-8b-rag.json before publishing.",
             size=9, color=TEXT_MUTED)


def slide_A2_confusion():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A2  ·  Confusion Matrices  ·  qwen3:8b + RAG")
    x, y, w, h = content_box()
    gap = Inches(0.4)
    cw = (w - gap) / 2

    def matrix(cx, cy, label, tp, fn, fp, tn):
        add_rect(s, cx, cy, cw, Inches(3.8), fill=CARD_BG,
                 line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.25), cy + Inches(0.15), cw - Inches(0.5),
                 Inches(0.4), label, size=14, bold=True, color=NAVY)
        # 2x2 grid
        gx = cx + Inches(0.6)
        gy = cy + Inches(0.95)
        cw2 = (cw - Inches(1.2)) / 2
        rh = Inches(1.0)
        # column headers
        add_text(s, gx, cy + Inches(0.6), cw2, Inches(0.3),
                 "Predicted: Vuln", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
        add_text(s, gx + cw2, cy + Inches(0.6), cw2, Inches(0.3),
                 "Predicted: Safe", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
        # row labels
        add_text(s, cx + Inches(0.05), gy, Inches(0.55), rh,
                 "Actual\nVuln", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.05), gy + rh, Inches(0.55), rh,
                 "Actual\nSafe", size=10, bold=True, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # cells
        cells = [(tp, GREEN, "TP"), (fn, RED, "FN"),
                 (fp, RED, "FP"), (tn, GREEN, "TN")]
        for i, (v, col, lbl) in enumerate(cells):
            r, c = i // 2, i % 2
            ccx = gx + c * cw2
            ccy = gy + r * rh
            bg = RGBColor(0xE9, 0xF3, 0xEA) if col == GREEN else RGBColor(0xF9, 0xE6, 0xEA)
            add_rect(s, ccx, ccy, cw2, rh, fill=bg, line=CARD_BORDER)
            add_text(s, ccx, ccy + Inches(0.05), cw2, Inches(0.4),
                     str(v), size=22, bold=True, color=NAVY_DARK,
                     align=PP_ALIGN.CENTER)
            add_text(s, ccx, ccy + Inches(0.55), cw2, Inches(0.3),
                     lbl, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    matrix(x, y, "Full corpus  (n = 101)", "50", "21", "8", "22")
    matrix(x + cw + gap, y, "Held-out test split  (n = 30)", "11", "7", "0", "12")
    add_text(s, x, y + Inches(3.95), w, Inches(0.3),
             "Replace cell counts with the actual canonical-taxonomy figures from evaluation/results/ before publishing.",
             size=9, color=TEXT_MUTED)


def slide_A3_manual_review():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A3  ·  Why the Single-Reviewer Manual Review Is Bounded")
    x, y, w, h = content_box()
    gap = Inches(0.25)
    cw = (w - gap) / 2
    ch = (Inches(3.6) - gap) / 2
    rubric = [
        ("Semantic correctness", "89.5 %", "Fix addresses the CWE."),
        ("Strategy alignment", "80 %", "Right repair pattern (parameterized query, escaping, etc.)."),
        ("Executability", "76 %", "Patch runs without modification."),
        ("Conciseness", "84 %", "Minimal diff, no scope creep."),
    ]
    for i, (lbl, val, sub) in enumerate(rubric):
        r, c = i // 2, i % 2
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        add_rect(s, cx, cy, cw, ch, fill=CARD_BG, line=CARD_BORDER, rounded=True)
        add_rect(s, cx, cy, cw, Inches(0.08), fill=NAVY)
        add_text(s, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4),
                 Inches(0.4), lbl, size=13, bold=True, color=NAVY_DARK)
        add_text(s, cx + Inches(0.2), cy + Inches(0.55), cw - Inches(0.4),
                 Inches(0.7), val, size=30, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.2), cy + Inches(1.3), cw - Inches(0.4),
                 ch - Inches(1.5), sub, size=10, color=TEXT_DARK)
    by = y + 2 * ch + gap + Inches(0.1)
    add_rect(s, x, by, w, Inches(0.7), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6), Inches(0.5),
             "Inter-rater on 10-sample subset = future work.  Headline R3 is the fully objective auto-applicable rate (90.32 %, n = 279).",
             size=12, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_A4_corpus():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A4  ·  Why a Curated Corpus, Not an Existing Benchmark")
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2

    # left — why others don't fit
    add_rect(s, x, y, cw, h - Inches(0.2), fill=CARD_BG,
             line=CARD_BORDER, rounded=True)
    add_text(s, x + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4), "WHY EXISTING JS/TS BENCHMARKS DON'T FIT", size=11,
             bold=True, color=ACCENT)
    items = [
        ("Juliet (NIST)", "C / C++ / Java only.  No first-party JS port."),
        ("OWASP Benchmark", "Java only."),
        ("SecurityEval (Pearce 2022)", "Small, partially synthetic.  CWE labelling inconsistent with OWASP Top 10."),
        ("CyberSecEval", "29.3 % of vulnerable samples are reproducible (Peng et al. 2025)."),
        ("CWEval (Peng 2025)", "119 curated tasks across 31 CWEs, 5 languages — peer-comparable to 101 cases."),
    ]
    ry = y + Inches(0.6)
    for name, sub in items:
        add_text(s, x + Inches(0.25), ry, cw - Inches(0.5), Inches(0.3),
                 name, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.25), ry + Inches(0.3), cw - Inches(0.5),
                 Inches(0.45), sub, size=10, color=TEXT_DARK)
        ry += Inches(0.75)

    # right — mitigations
    rx = x + cw + gap
    add_rect(s, rx, y, cw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4),
             "SELECTION-BIAS MITIGATIONS IN THIS THESIS", size=11, bold=True,
             color=ACCENT)
    mits = [
        ("Held-out 71 / 30 split", "Frozen TEST set keeps threshold tuning out of headline numbers."),
        ("15-case external corpus", "NodeGoat, Juice Shop, 3 named CVEs — independent of the curator."),
        ("Whole-project NodeGoat run", "Project-level recall independent of sample selection."),
        ("Provenance bound into manifest", "Ed25519 signature; corpus cannot be silently rewritten."),
    ]
    ry = y + Inches(0.65)
    for name, sub in mits:
        add_text(s, rx + Inches(0.25), ry, cw - Inches(0.5), Inches(0.3),
                 name, size=12, bold=True, color=TEXT_LIGHT)
        add_text(s, rx + Inches(0.25), ry + Inches(0.3), cw - Inches(0.5),
                 Inches(0.55), sub, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
        ry += Inches(0.95)


def slide_A5_no_cloud():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A5  ·  Why No Cloud-LLM Baseline")
    x, y, w, h = content_box()
    gap = Inches(0.3)
    cw = (w - gap) / 2

    # left — threat model diagram (placeholder)
    placeholder(s, x, y, cw, h - Inches(0.2),
                "THREAT-MODEL DIAGRAM",
                sub="Laptop  →  loopback boundary (dashed)  ✗  CLOUD baseline outside boundary  ·  ✅  SAST + Code Guardian inside")

    # right — argument
    rx = x + cw + gap
    add_rect(s, rx, y, cw, h - Inches(0.2), fill=NAVY, rounded=True)
    add_text(s, rx + Inches(0.25), y + Inches(0.15), cw - Inches(0.5),
             Inches(0.4),
             "ARGUMENT", size=11, bold=True, color=ACCENT)
    pts = [
        ("Cloud-LLM is the problem statement",
         "Including it as a baseline contradicts the threat model."),
        ("Floor provided by 3 SAST baselines",
         "Semgrep, CodeQL, ESLint  ·  all local, all reproducible."),
        ("Frontier cloud will always win raw F1",
         "The thesis competes on the constraint set:  local, zero-egress, reproducible."),
        ("Post-defense addendum",
         "A one-paragraph comparison to published GPT-4 numbers on similar JS corpora is straightforward."),
    ]
    ry = y + Inches(0.6)
    for name, sub in pts:
        add_text(s, rx + Inches(0.25), ry, cw - Inches(0.5), Inches(0.35),
                 name, size=12, bold=True, color=TEXT_LIGHT)
        add_text(s, rx + Inches(0.25), ry + Inches(0.35), cw - Inches(0.5),
                 Inches(0.7), sub, size=10, color=RGBColor(0xC8, 0xD0, 0xE0))
        ry += Inches(1.05)


def slide_A6_consensus():
    s = prs.slides.add_slide(BLANK)
    appendix_header(s, "A6  ·  Why R2 Reports 100 % Inter-Run Agreement")
    x, y, w, h = content_box()
    gap = Inches(0.2)
    cw = (w - 2 * gap) / 3

    # three identical JSON blobs
    json_text = ('{\n  "findings": [\n    {\n      "category":\n        "sql-injection",\n'
                 '      "severity":\n        "high",\n      "line": 42\n    }\n  ]\n}')
    for i in range(3):
        cx = x + i * (cw + gap)
        add_rect(s, cx, y, cw, Inches(3.5),
                 fill=NAVY_DARK, line=CARD_BORDER, rounded=True)
        add_text(s, cx + Inches(0.2), y + Inches(0.1), cw - Inches(0.4),
                 Inches(0.35), f"RUN {i + 1}", size=10, bold=True, color=ACCENT)
        add_text(s, cx + Inches(0.2), y + Inches(0.45), cw - Inches(0.4),
                 Inches(3.0), json_text, size=11, color=TEXT_LIGHT,
                 font="Courier New")

    # bottom — explanation
    by = y + Inches(3.7)
    add_rect(s, x, by, w, Inches(0.95), fill=NAVY, rounded=True)
    add_text(s, x + Inches(0.3), by + Inches(0.1), w - Inches(0.6),
             Inches(0.3),
             "WHY THEY'RE IDENTICAL", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), by + Inches(0.4), w - Inches(0.6),
             Inches(0.5),
             "temperature = 0   ·   seed = 42   ·   format = JSON   ·   fixed prompt.  Byte-identical → consensus signal requires seed rotation (future work).",
             size=12, color=TEXT_LIGHT)


# ============================================================
# BUILD
# ============================================================
slide_01_title()
slide_02_problem_illustration()
slide_03_two_failures()
slide_04_three_gap()
slide_05_rqs_reqs()
slide_06_existing_approaches()
slide_07_gap_matrix()
slide_08_workflows()
slide_09_two_stage()
slide_10_screenshot()
slide_11_architecture()
slide_12_privacy_cards()
slide_13_eval_setup()
slide_14_hero()
slide_15_r1_accuracy()
slide_16_rag_ablation()
slide_17_r3_repair()
slide_18_r4_latency()
slide_19_r5_privacy()
slide_20_takeaways()
slide_21_limitations()
slide_22_future_work()
slide_23_contributions()
slide_24_thanks()
slide_A1_full_results()
slide_A2_confusion()
slide_A3_manual_review()
slide_A4_corpus()
slide_A5_no_cloud()
slide_A6_consensus()

# Replace placeholder "__" page count
total = 24
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if "/ __" in r.text:
                        r.text = r.text.replace("/ __", f"/ {total}")

out = "/Users/hafiz/personal/repos/tuc-master-thesis/thesis-presentation/defense_slides_v2.pptx"
prs.save(out)
print(f"OK  wrote {out}  ·  {len(prs.slides)} slides")
